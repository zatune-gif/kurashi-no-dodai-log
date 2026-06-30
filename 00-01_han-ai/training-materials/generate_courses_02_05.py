"""
コース②〜⑤ スライド一括生成スクリプト
Google Slides にインポートして使用する .pptx を生成します。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── カラーパレット ──────────────────────────────────────
C_PRIMARY   = RGBColor(0x1E, 0x40, 0xAF)
C_ACCENT    = RGBColor(0xDA, 0x77, 0x56)
C_TEXT      = RGBColor(0x1E, 0x29, 0x3B)
C_MUTED     = RGBColor(0x64, 0x74, 0x8B)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_LIGHT  = RGBColor(0xEF, 0xF6, 0xFF)
C_TABLE_ALT = RGBColor(0xF8, 0xFA, 0xFC)
C_GREEN     = RGBColor(0x05, 0x96, 0x69)
C_RED       = RGBColor(0xDC, 0x26, 0x26)
C_BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
C_LIGHT_BLUE = RGBColor(0xBF, 0xDB, 0xFE)
C_PALE_BLUE  = RGBColor(0x93, 0xC5, 0xFD)
C_GREEN_BG  = RGBColor(0xEC, 0xFD, 0xF5)
C_AMBER     = RGBColor(0xD9, 0x77, 0x06)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ── 共通ヘルパー ─────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, width, height, fill_color):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def tb(slide, text, left, top, width, height,
       size=14, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb

def add_line(tf, text, size=13, bold=False, italic=False, color=None):
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color

def add_blank_line(tf, size=6):
    p = tf.add_paragraph()
    p.add_run().font.size = Pt(size)

def add_table(slide, headers, rows, left, top, width, height,
              header_bg=None, col_widths=None, alt_row=True,
              h_size=12, b_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, ratio in enumerate(col_widths):
            tbl.columns[i].width = int(width * ratio)
    else:
        cw = width // n_cols
        for i in range(n_cols):
            tbl.columns[i].width = cw

    hbg = header_bg or C_PRIMARY

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hbg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(h_size)
        run.font.color.rgb = C_WHITE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if alt_row and ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TABLE_ALT
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(b_size)
            run.font.color.rgb = C_TEXT

    return tbl

def slide_header(slide, title, time_str=None):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.65), C_PRIMARY)
    tb(slide, title, Inches(0.3), Inches(0.07), Inches(8.5), Inches(0.55),
       size=18, bold=True, color=C_WHITE)
    if time_str:
        tb(slide, f"[{time_str}]", Inches(7.8), Inches(0.1), Inches(2.0), Inches(0.5),
           size=12, color=C_PALE_BLUE, align=PP_ALIGN.RIGHT)

def info_box(slide, text, left, top, width, height, bg=None, text_color=None, size=14, bold=False):
    rect(slide, left, top, width, height, bg or C_BG_LIGHT)
    tb(slide, text, left + Inches(0.15), top + Inches(0.07),
       width - Inches(0.3), height - Inches(0.1),
       size=size, bold=bold, color=text_color or C_PRIMARY)

def slide_section(prs, num, title, time_str, description="", color=None):
    sl = blank_slide(prs)
    set_bg(sl, C_PRIMARY)
    badge_color = color or C_ACCENT
    badge = rect(sl, Inches(0.6), Inches(0.9), Inches(1.2), Inches(1.2), badge_color)
    tb(sl, str(num), Inches(0.6), Inches(0.95), Inches(1.2), Inches(1.1),
       size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, title, Inches(2.2), Inches(1.05), Inches(7.5), Inches(1.0),
       size=30, bold=True, color=C_WHITE)
    tb(sl, f"[{time_str}]", Inches(2.2), Inches(2.1), Inches(5), Inches(0.5),
       size=14, color=C_PALE_BLUE)
    if description:
        tb(sl, description, Inches(2.2), Inches(2.7), Inches(7.5), Inches(1.0),
           size=15, color=C_LIGHT_BLUE)

def slide_break(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_PRIMARY)
    tb(sl, "- 休憩 -", Inches(0), Inches(2.0), SLIDE_W, Inches(1.0),
       size=40, bold=True, color=C_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    tb(sl, "5〜10分", Inches(0), Inches(3.1), SLIDE_W, Inches(0.6),
       size=18, color=C_PALE_BLUE, align=PP_ALIGN.CENTER)

def slide_common_outro(prs, tool_hint=""):
    sl = blank_slide(prs)
    slide_header(sl, "まとめ・サービス案内")
    rect(sl, Inches(0.5), Inches(0.85), Inches(9), Inches(1.8), C_BG_LIGHT)
    tb(sl, "次のステップ | ざつね屋のサービス案内",
       Inches(0.7), Inches(0.95), Inches(8.5), Inches(0.45),
       size=13, bold=True, color=C_PRIMARY)
    svc = slide.shapes.add_textbox if False else sl.shapes.add_textbox(
        Inches(0.7), Inches(1.42), Inches(8.5), Inches(1.1))
    tf = svc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "📚 他の実践編コース（文書・画像・動画・Claude Code）で幅をさらに広げる"
    r.font.size = Pt(12)
    r.font.color.rgb = C_TEXT
    add_line(tf, "🛠 AI業務改善オーダーメイド・🤝 AI開発伴走 — 自社の課題を一緒に解決する",
             size=12, color=C_TEXT)
    if tool_hint:
        add_line(tf, tool_hint, size=11, italic=True, color=C_MUTED)

    tb(sl, "本日の学びをまとめましょう",
       Inches(0.5), Inches(2.85), Inches(9), Inches(0.45),
       size=14, bold=True, color=C_TEXT)
    tb(sl, "「わからないことはAIに聞いてみてください」",
       Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
       size=15, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    tb(sl, "質問・ご相談はお気軽にどうぞ",
       Inches(0.5), Inches(4.0), Inches(9), Inches(0.5),
       size=13, color=C_MUTED, align=PP_ALIGN.CENTER)
    return sl


# ═══════════════════════════════════════════════════════════
#  コース②：実践編（文書系）
# ═══════════════════════════════════════════════════════════

def make_course02(output_path):
    prs = new_prs()

    # ── タイトル ──
    sl = blank_slide(prs)
    set_bg(sl, C_PRIMARY)
    rect(sl, Inches(0), Inches(3.7), SLIDE_W, Inches(0.06), C_ACCENT)
    tb(sl, "ざつね屋", Inches(0.6), Inches(0.4), Inches(4), Inches(0.45),
       size=13, color=C_PALE_BLUE)
    tb(sl, "コース ②", Inches(0.6), Inches(0.95), Inches(9), Inches(0.55),
       size=18, color=C_LIGHT_BLUE)
    tb(sl, "実践編（文書系）", Inches(0.6), Inches(1.5), Inches(9), Inches(1.3),
       size=44, bold=True, color=C_WHITE)
    tb(sl, "AIで文書作業を大幅効率化する 120分",
       Inches(0.6), Inches(2.85), Inches(9), Inches(0.6),
       size=16, color=C_LIGHT_BLUE)
    tb(sl, "ツール：Gemini（無料版）　｜　成果物：業務別プロンプトライブラリ10本",
       Inches(0.6), Inches(4.1), Inches(9), Inches(0.5),
       size=13, color=C_PALE_BLUE)

    # ── アジェンダ ──
    sl = blank_slide(prs)
    slide_header(sl, "本日のアジェンダ")
    add_table(sl, ["#", "セクション", "時間"],
        [["1","イントロ","10分"],
         ["2","社内文書（議事録・メール・報告書）","30分"],
         ["3","広報・発信文書（SNS・チラシ・HP）","30分"],
         ["—","休憩","10分"],
         ["4","採用広報（求人票・会社紹介文）","30分"],
         ["5","ライブラリ化・保存","5分"],
         ["6","まとめ","5分"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.6),
        col_widths=[0.06, 0.72, 0.22])

    # ── Section 1：イントロ ──
    slide_section(prs, "1", "イントロ", "10分", "Gemini確認 ＆ 今日のゴール説明")
    sl = blank_slide(prs)
    slide_header(sl, "イントロ — 今日のゴール＆ツール確認", "10分")
    add_table(sl, ["確認項目", "内容"],
        [["Geminiログイン","Googleアカウントでgemini.google.comへアクセス（未作成の人はその場で作成）"],
         ["コース①未受講の方","プロンプト4要素を3分でおさらい → 「役割・背景・指示・形式」"],
         ["今日のゴール","文書系3カテゴリ（社内・広報・採用）でプロンプトを作り、10本ライブラリ化して持ち帰る"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.5),
        col_widths=[0.28, 0.72])
    info_box(sl, "ひな型（全セクション共通）：「あなたは[役割]です。[状況]の[文書の種類]を[形式]で作成してください」",
             Inches(0.5), Inches(3.55), Inches(9), Inches(0.65),
             text_color=C_TEXT, size=13)

    # ── Section 2：社内文書 ──
    slide_section(prs, "2", "社内文書", "30分", "議事録・メール・報告書・マニュアルを作る")
    sl = blank_slide(prs)
    slide_header(sl, "社内文書 — 実践①（全員共通）：議事録", "実践15分")
    info_box(sl, "お題：「先週の会議の発言メモを渡すから議事録に整形して」",
             Inches(0.5), Inches(0.85), Inches(9), Inches(0.65),
             bg=C_BG_LIGHT, text_color=C_TEXT, size=14, bold=True)
    add_table(sl, ["Step", "内容", "時間"],
        [["1","発言メモを手元に準備（なければ架空で可）","2分"],
         ["2","ひな型を使ってプロンプトを入力 → 送信","5分"],
         ["3","出力を確認・修正指示で調整（「もっと箇条書きに」等）","5分"],
         ["4","見せ合い・感想共有","3分"]],
        Inches(0.5), Inches(1.68), Inches(9), Inches(2.6),
        col_widths=[0.07, 0.73, 0.20])
    tb(sl, "ポイント：AIへの入力は「コピペ」でOK。整形・要約は得意分野です",
       Inches(0.5), Inches(4.45), Inches(9), Inches(0.6),
       size=13, bold=True, color=C_PRIMARY)

    sl = blank_slide(prs)
    slide_header(sl, "社内文書 — 実践②（自由テーマ）：想起リスト", "実践10分 + 見せ合い5分")
    tb(sl, "こんな文書もAIで作れます — 好きなテーマを選んでプロンプトを試してみてください",
       Inches(0.5), Inches(0.85), Inches(9), Inches(0.55),
       size=13, color=C_MUTED)
    add_table(sl, ["テーマ", "使い道・例"],
        [["メール返信文","クレーム対応・お礼・依頼・問い合わせへの返信"],
         ["報告書のたたき台","月次・週次・プロジェクト報告書の素材作り"],
         ["社内マニュアル","業務手順・ルール説明をわかりやすくまとめる"],
         ["社内向けお知らせ文","休業告知・システム変更・ルール変更のアナウンス"],
         ["上司への相談文","困っていること・提案・改善要望を整理して伝える"]],
        Inches(0.5), Inches(1.5), Inches(9), Inches(3.55),
        col_widths=[0.30, 0.70])

    # ── Section 3：広報・発信文書 ──
    slide_section(prs, "3", "広報・発信文書", "30分", "SNS投稿・チラシ・HPテキストを作る")
    sl = blank_slide(prs)
    slide_header(sl, "広報・発信文書 — 実践①（全員共通）：SNS投稿文", "実践15分")
    info_box(sl, "お題：「自社の商品・サービス・お知らせを Instagram / Facebook 投稿1本分作って」",
             Inches(0.5), Inches(0.85), Inches(9), Inches(0.65),
             bg=C_BG_LIGHT, text_color=C_TEXT, size=14, bold=True)
    add_table(sl, ["Step", "内容", "時間"],
        [["1","自社の投稿テーマを決める（商品紹介・イベント・お知らせ等）","2分"],
         ["2","プロンプト入力：「会社名・業種・投稿テーマ・ハッシュタグの有無」を指定","5分"],
         ["3","出力確認 → 「もっとカジュアルに」「絵文字を増やして」等で調整","5分"],
         ["4","見せ合い・感想共有","3分"]],
        Inches(0.5), Inches(1.68), Inches(9), Inches(2.6),
        col_widths=[0.07, 0.73, 0.20])
    tb(sl, "ポイント：AIの文章はそのまま使わず「人間のチェック」を入れることで質が上がる",
       Inches(0.5), Inches(4.45), Inches(9), Inches(0.6),
       size=13, bold=True, color=C_PRIMARY)

    sl = blank_slide(prs)
    slide_header(sl, "広報・発信文書 — 実践②（自由テーマ）：想起リスト", "実践10分 + 見せ合い5分")
    tb(sl, "こんな発信もAIで下書きできます",
       Inches(0.5), Inches(0.85), Inches(9), Inches(0.55),
       size=13, color=C_MUTED)
    add_table(sl, ["テーマ", "使い道・例"],
        [["チラシ・DM文案","イベント告知・セール・新商品紹介チラシのコピー"],
         ["HP更新テキスト","会社紹介・サービス説明・ニュースページの文章"],
         ["メルマガ・ニュースレター","定期配信の構成と文面をたたき台で作る"],
         ["プレスリリース","新サービス・受賞・採用募集などの告知文"],
         ["キャッチコピー案","商品・サービスの魅力を一言で伝えるフレーズを複数案出す"]],
        Inches(0.5), Inches(1.5), Inches(9), Inches(3.55),
        col_widths=[0.30, 0.70])

    # ── 休憩 ──
    slide_break(prs)

    # ── Section 4：採用広報 ──
    slide_section(prs, "4", "採用広報", "30分", "求人票・会社紹介文・採用SNS投稿を作る")
    sl = blank_slide(prs)
    slide_header(sl, "採用広報 — 実践①（全員共通）：求人票", "実践15分")
    info_box(sl, "お題：「職種・条件・自社の魅力を入れて求人票を作って」",
             Inches(0.5), Inches(0.85), Inches(9), Inches(0.65),
             bg=C_BG_LIGHT, text_color=C_TEXT, size=14, bold=True)
    add_table(sl, ["Step", "内容", "時間"],
        [["1","募集職種と最低限の条件（給与・勤務地・仕事内容）を手元に用意","2分"],
         ["2","「未経験歓迎・自社の強み・応募しやすいポイント」も含めてプロンプト入力","5分"],
         ["3","出力確認 → 「もう少し求職者の不安を解消する言い方に」等で調整","5分"],
         ["4","見せ合い・感想共有","3分"]],
        Inches(0.5), Inches(1.68), Inches(9), Inches(2.6),
        col_widths=[0.07, 0.73, 0.20])
    tb(sl, "ポイント：自社の強みを「普通の言葉」でAIに伝えると、うまく文章化してくれます",
       Inches(0.5), Inches(4.45), Inches(9), Inches(0.6),
       size=13, bold=True, color=C_PRIMARY)

    sl = blank_slide(prs)
    slide_header(sl, "採用広報 — 実践②（自由テーマ）：想起リスト", "実践10分 + 見せ合い5分")
    tb(sl, "採用活動全体をAIがサポートします",
       Inches(0.5), Inches(0.85), Inches(9), Inches(0.55),
       size=13, color=C_MUTED)
    add_table(sl, ["テーマ", "使い道・例"],
        [["会社紹介文","採用ページ・合同説明会・会社案内の文章を作る"],
         ["採用広報SNS投稿","社員の一日・職場環境・理念を発信する投稿文"],
         ["面接でよく聞く質問リスト","職種別の想定質問と評価ポイントをまとめる"],
         ["内定通知文","内定のご連絡メール・書面のたたき台"],
         ["入社前案内文","初日の持ち物・手続き・注意事項をまとめた案内文"]],
        Inches(0.5), Inches(1.5), Inches(9), Inches(3.55),
        col_widths=[0.30, 0.70])

    # ── Section 5：ライブラリ化 ──
    slide_section(prs, "5", "ライブラリ化・保存", "5分", "今日作ったプロンプトを「明日から使えるテンプレート」に仕上げる")
    sl = blank_slide(prs)
    slide_header(sl, "プロンプトライブラリ化 — テンプレートに変換して保存", "5分")
    add_table(sl, ["Step", "内容"],
        [["1 振り返り","セクション2〜4で作ったプロンプトを見直し・気に入った10本を確認"],
         ["2 テンプレート化","固有情報を [ ] に置き換えて使い回せる形にする"],
         ["3 保存","Googleドキュメント or メモアプリにコピーして持ち帰り"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.2),
        col_widths=[0.25, 0.75])
    tb(sl, "テンプレート化の例：",
       Inches(0.5), Inches(3.2), Inches(9), Inches(0.4),
       size=13, bold=True, color=C_TEXT)
    before_box = sl.shapes.add_textbox(Inches(0.5), Inches(3.65), Inches(9), Inches(0.55))
    tf = before_box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "変換前：「広島県の製造業・従業員30名の会社で未経験歓迎の溶接工を募集」"
    r.font.size = Pt(12)
    r.font.color.rgb = C_MUTED
    add_line(tf, "変換後：「[都道府県]の[業種]・従業員[人数]名の会社で[条件]の[職種]を募集」",
             size=13, bold=True, color=C_GREEN)

    # ── Section 6：まとめ ──
    slide_section(prs, "6", "まとめ", "5分")
    sl = blank_slide(prs)
    slide_header(sl, "まとめ")
    rect(sl, Inches(0.5), Inches(0.85), Inches(9), Inches(1.75), C_BG_LIGHT)
    tb(sl, "次のステップ | ざつね屋のサービス案内",
       Inches(0.7), Inches(0.95), Inches(8.5), Inches(0.45),
       size=13, bold=True, color=C_PRIMARY)
    svc = sl.shapes.add_textbox(Inches(0.7), Inches(1.42), Inches(8.5), Inches(1.0))
    tf = svc.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "📚 実践編（画像・動画・Claude Code）— 発信力・自動化をさらに強化する"
    r.font.size = Pt(12)
    r.font.color.rgb = C_TEXT
    add_line(tf, "🛠 AI業務改善オーダーメイド・🤝 AI開発伴走 — 自社の課題を一緒に解決する", size=12, color=C_TEXT)
    tb(sl, "「わからないことはAIに聞いてみてください」",
       Inches(0.5), Inches(3.0), Inches(9), Inches(0.6),
       size=16, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    tb(sl, "本日はありがとうございました",
       Inches(0.5), Inches(3.8), Inches(9), Inches(0.5),
       size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    print(f"[OK] course02: {output_path} ({len(prs.slides)} slides)")


# ═══════════════════════════════════════════════════════════
#  コース③：実践編（画像系）
# ═══════════════════════════════════════════════════════════

def make_course03(output_path):
    prs = new_prs()

    # ── タイトル ──
    sl = blank_slide(prs)
    set_bg(sl, C_PRIMARY)
    rect(sl, Inches(0), Inches(3.7), SLIDE_W, Inches(0.06), C_ACCENT)
    tb(sl, "ざつね屋", Inches(0.6), Inches(0.4), Inches(4), Inches(0.45),
       size=13, color=C_PALE_BLUE)
    tb(sl, "コース ③", Inches(0.6), Inches(0.95), Inches(9), Inches(0.55),
       size=18, color=C_LIGHT_BLUE)
    tb(sl, "実践編（画像系）", Inches(0.6), Inches(1.5), Inches(9), Inches(1.3),
       size=44, bold=True, color=C_WHITE)
    tb(sl, "デザイナー不在でもSNS・チラシ・広報のビジュアルが作れる 120分",
       Inches(0.6), Inches(2.85), Inches(9), Inches(0.6),
       size=15, color=C_LIGHT_BLUE)
    tb(sl, "ツール：Canva 無料版のみ　｜　成果物：SNS画像 + AI生成画像 + チラシ 各1点",
       Inches(0.6), Inches(4.1), Inches(9), Inches(0.5),
       size=13, color=C_PALE_BLUE)

    # ── アジェンダ ──
    sl = blank_slide(prs)
    slide_header(sl, "本日のアジェンダ")
    add_table(sl, ["#", "セクション", "内容", "時間"],
        [["1","イントロ","Canvaログイン確認・基本画面の説明","10分"],
         ["2","SNS投稿画像","Canvaテンプレートで編集実践","30分"],
         ["3","AI画像生成","Magic Mediaでプロンプト入力→生成","30分"],
         ["—","休憩","","10分"],
         ["4","チラシ作成","テンプレ＋AI画像を組み合わせて仕上げる","30分"],
         ["5","まとめ","ダウンロード・共有・活用アドバイス","10分"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.55),
        col_widths=[0.06, 0.22, 0.52, 0.20])

    # ── Section 1：イントロ ──
    slide_section(prs, "1", "イントロ", "10分", "Canvaにログインして基本画面を確認")
    sl = blank_slide(prs)
    slide_header(sl, "イントロ — Canva基本画面の説明", "10分")
    add_table(sl, ["確認項目", "内容"],
        [["Canvaログイン","canva.com でGoogleアカウント or メールアドレスで登録・ログイン（無料版）"],
         ["新規デザイン作成","「デザインを作成」→ 用途を選ぶ（SNS投稿・チラシ等）"],
         ["テンプレート","Canvaが用意している「使える雛形」を選んで文字・色・画像を差し替える"],
         ["Magic Media","AIで画像生成できる機能（無料枠：月50回程度）"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.8),
        col_widths=[0.25, 0.75])
    tb(sl, "今日の流れ：Canva操作を覚える → AI素材を作る → 両方を組み合わせて完成品を作る",
       Inches(0.5), Inches(3.85), Inches(9), Inches(0.6),
       size=14, bold=True, color=C_PRIMARY)

    # ── テンプレートの意味の違い ──
    sl = blank_slide(prs)
    slide_header(sl, "「テンプレート」の意味の違い — 本日の重要ポイント")
    add_table(sl, ["セクション", "「テンプレ」の意味"],
        [["2・4（SNS画像・チラシ）","Canvaのデザインテンプレート ＝ Canvaが用意している雛形を選んで編集する"],
         ["3（AI画像生成）","プロンプトのひな型 ＝ 「[スタイル]の[主題]、[用途]向け」のような[ ]を埋めてMagic Mediaに入力する形式"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.3),
        col_widths=[0.30, 0.70])
    info_box(sl, "混乱しやすいポイントなので、都度「こちらはデザインの雛形です / こちらはAIへの指示文です」と明示してください",
             Inches(0.5), Inches(3.35), Inches(9), Inches(0.75),
             bg=RGBColor(0xFF, 0xF0, 0xE8), text_color=C_ACCENT, size=13)

    # ── Section 2：SNS投稿画像 ──
    slide_section(prs, "2", "SNS投稿画像", "30分", "Canvaのデザインテンプレートで実践")
    sl = blank_slide(prs)
    slide_header(sl, "SNS投稿画像 — 実践①（全員共通）：Instagram投稿", "実践15分")
    info_box(sl, "お題：自社のサービス・商品紹介の Instagram 投稿画像を1枚作る",
             Inches(0.5), Inches(0.85), Inches(9), Inches(0.65),
             text_color=C_TEXT, size=14, bold=True)
    add_table(sl, ["Step", "内容", "時間"],
        [["1","「デザインを作成」→「Instagramの投稿（正方形）」を選択","2分"],
         ["2","テンプレートを検索・選択（業種・雰囲気で選ぶ）","3分"],
         ["3","テキストを自社の内容に書き換え・色やフォントを調整","7分"],
         ["4","ダウンロード（PNG）して確認","3分"]],
        Inches(0.5), Inches(1.68), Inches(9), Inches(2.4),
        col_widths=[0.07, 0.73, 0.20])
    tb(sl, "ポイント：テンプレートのレイアウトはそのままでOK。テキストと色だけ変えるだけで「自社らしい」デザインになります",
       Inches(0.5), Inches(4.25), Inches(9), Inches(0.7),
       size=13, bold=True, color=C_PRIMARY)

    sl = blank_slide(prs)
    slide_header(sl, "SNS投稿画像 — 実践②（自由テーマ）：想起リスト", "実践10分")
    tb(sl, "こんなビジュアルもCanvaで作れます",
       Inches(0.5), Inches(0.85), Inches(9), Inches(0.5),
       size=13, color=C_MUTED)
    add_table(sl, ["テーマ", "作り方・ポイント"],
        [["Facebook投稿画像","横長（1200×628px）でテンプレートを選択"],
         ["キャンペーン・セール告知","「セール」「SALE」で検索するとバナー向きテンプレが出る"],
         ["季節のご挨拶画像","「年末年始」「春」「夏」等で検索・差し替えるだけでOK"],
         ["お客様の声・事例紹介","引用デザインのテンプレートを選んでテキスト差し替え"],
         ["LINE公式アカウント用バナー","1200×600pxで作成・ダウンロードしてLINE管理画面でアップ"]],
        Inches(0.5), Inches(1.45), Inches(9), Inches(3.85),
        col_widths=[0.30, 0.70])

    # ── Section 3：AI画像生成 ──
    slide_section(prs, "3", "AI画像生成", "30分", "Magic Mediaでテキストから画像を生成する")
    sl = blank_slide(prs)
    slide_header(sl, "AI画像生成 — Canva Magic Media の使い方", "実践25分")
    add_table(sl, ["操作手順", "内容"],
        [["1 Magic Mediaを開く","Canvaエディタ左パネル「アプリ」→「Magic Media」を選択"],
         ["2 プロンプトを入力","テキスト欄にひな型を使って指示を入力"],
         ["3 スタイルを選ぶ","写真風・イラスト・水彩・3D等から選択（または「なし」でも可）"],
         ["4 生成 → 選択","4枚生成される → 気に入ったものを選んでデザインに追加"],
         ["5 組み合わせ","CanvaのテンプレートレイアウトにAI生成画像を貼り付けて仕上げる"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.4),
        col_widths=[0.28, 0.72])
    info_box(sl, "無料枠：月50回程度（目安）。生成が「上限」と表示されたら講師に相談してください",
             Inches(0.5), Inches(4.43), Inches(9), Inches(0.6),
             bg=RGBColor(0xFF, 0xF0, 0xE8), text_color=C_ACCENT, size=12)

    sl = blank_slide(prs)
    slide_header(sl, "AI画像生成 — プロンプトひな型＆想起リスト")
    info_box(sl, "ひな型：「[スタイル]で[主題・テーマ]の画像、[用途]向け、[雰囲気]な印象で」",
             Inches(0.5), Inches(0.85), Inches(9), Inches(0.65),
             text_color=C_TEXT, size=14, bold=True)
    tb(sl, "例：「水彩画風で春の桜と笑顔のスタッフの画像、SNS投稿向け、明るく温かみのある印象で」",
       Inches(0.5), Inches(1.62), Inches(9), Inches(0.5),
       size=13, italic=True, color=C_MUTED)
    add_table(sl, ["生成したいもの", "プロンプトのポイント"],
        [["会社・店舗の雰囲気画像","業種・場所・雰囲気（清潔感・温かみ・プロフェッショナル等）を指定"],
         ["商品・サービスのイメージ","商品の特徴・使用シーン・ターゲット感を伝える"],
         ["SNS用の背景・素材","「シンプル背景」「テキスト入りやすい余白あり」等を指定"],
         ["採用広報用のシーン画像","「職場で働くスタッフ」「チームワーク」等・人物は指定しにくい場合あり"]],
        Inches(0.5), Inches(2.3), Inches(9), Inches(2.85),
        col_widths=[0.35, 0.65])

    # ── 休憩 ──
    slide_break(prs)

    # ── Section 4：チラシ作成 ──
    slide_section(prs, "4", "チラシ作成", "30分", "テンプレート＋AI画像を組み合わせて仕上げる")
    sl = blank_slide(prs)
    slide_header(sl, "チラシ作成 — セクション2・3の技術を組み合わせる", "30分")
    add_table(sl, ["Step", "内容", "時間"],
        [["1","「デザインを作成」→「A4チラシ」（または任意のサイズ）を選択","3分"],
         ["2","テンプレートを選択（業種・用途で絞り込み）","5分"],
         ["3","セクション3で生成したAI画像をアップロード or 再生成して挿入","7分"],
         ["4","テキスト（キャッチコピー・住所・電話番号等）を自社の内容に差し替え","10分"],
         ["5","PDF または PNG でダウンロード","5分"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.65),
        col_widths=[0.07, 0.73, 0.20])
    tb(sl, "完成しなくても大丈夫。続きはCanvaアカウントにログインすれば、いつでも再開できます",
       Inches(0.5), Inches(4.7), Inches(9), Inches(0.6),
       size=13, italic=True, color=C_MUTED)

    # ── Section 5：まとめ ──
    slide_section(prs, "5", "まとめ", "10分", "ダウンロード・共有・活用アドバイス")
    sl = blank_slide(prs)
    slide_header(sl, "まとめ — 活用アドバイス＆サービス案内", "10分")
    add_table(sl, ["活用ポイント", "内容"],
        [["ダウンロード","PNG（SNS・メール用）・PDF（印刷用）どちらでも書き出せる"],
         ["後から編集","Canvaアカウントにログインすれば続きから編集・複製して使い回せる"],
         ["社内共有","共有リンクを送るだけで他のメンバーも閲覧・編集できる"],
         ["ブランドの統一","ブランドカラー・フォントを「ブランド設定」に登録して毎回使う"],
         ["テンプレとして使い回す","一度作ったデザインを複製して内容だけ変えれば量産可能"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.3),
        col_widths=[0.25, 0.75])
    rect(sl, Inches(0.5), Inches(4.35), Inches(9), Inches(0.65), C_BG_LIGHT)
    tb(sl, "次のステップ：実践編（動画系）でCanvaの動画機能、または実践編（文書系）でテキスト発信の質も上げましょう",
       Inches(0.7), Inches(4.42), Inches(8.5), Inches(0.55),
       size=13, bold=True, color=C_PRIMARY)

    prs.save(output_path)
    print(f"[OK] course03: {output_path} ({len(prs.slides)} slides)")


# ═══════════════════════════════════════════════════════════
#  コース④：実践編（動画系）
# ═══════════════════════════════════════════════════════════

def make_course04(output_path):
    prs = new_prs()

    # ── タイトル ──
    sl = blank_slide(prs)
    set_bg(sl, C_PRIMARY)
    rect(sl, Inches(0), Inches(3.7), SLIDE_W, Inches(0.06), C_ACCENT)
    tb(sl, "ざつね屋", Inches(0.6), Inches(0.4), Inches(4), Inches(0.45),
       size=13, color=C_PALE_BLUE)
    tb(sl, "コース ④", Inches(0.6), Inches(0.95), Inches(9), Inches(0.55),
       size=18, color=C_LIGHT_BLUE)
    tb(sl, "実践編（動画系）", Inches(0.6), Inches(1.5), Inches(9), Inches(1.3),
       size=44, bold=True, color=C_WHITE)
    tb(sl, "複数のAIが連携して1本の動画が完成する体験 120分",
       Inches(0.6), Inches(2.85), Inches(9), Inches(0.6),
       size=16, color=C_LIGHT_BLUE)
    tb(sl, "ツール：CapCut 無料版 ＋ Gemini（台本）　｜　成果物：ショート動画1本",
       Inches(0.6), Inches(4.1), Inches(9), Inches(0.5),
       size=13, color=C_PALE_BLUE)

    # ── アジェンダ ──
    sl = blank_slide(prs)
    slide_header(sl, "本日のアジェンダ")
    add_table(sl, ["#", "セクション", "内容", "時間"],
        [["1","イントロ","CapCutログイン確認・基本画面の説明","10分"],
         ["2","AI動画生成","Seedanceでテキスト→映像生成","15分"],
         ["3","ナレーション追加","Geminiで台本生成 → CapCut TTSで読み上げ","15分"],
         ["4","AI自動字幕","字幕生成（確認まで）","10分"],
         ["—","休憩","","10分"],
         ["5","基本編集・書き出し","字幕修正＋クリップ編集＋書き出し","50分"],
         ["6","まとめ","ダウンロード・SNS投稿・活用アドバイス","10分"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.55),
        col_widths=[0.06, 0.22, 0.52, 0.20])

    # ── ストーリーライン ──
    sl = blank_slide(prs)
    slide_header(sl, "今日のコンセプト — 複数のAIが連携して1本の動画を作る")
    flow = sl.shapes.add_textbox(Inches(0.5), Inches(0.88), Inches(9), Inches(1.0))
    tf = flow.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "映像AI（Seedance）→ 文書AI（Gemini 台本）→ 読み上げAI（CapCut TTS）→ 字幕AI（CapCut）→ 人間が最終編集"
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = C_PRIMARY
    add_table(sl, ["担当AI", "役割", "ツール"],
        [["映像AI","テキスト指示から映像クリップを生成","CapCut Seedance"],
         ["文書AI","ナレーション原稿（台本）を生成","Gemini（無料版）"],
         ["読み上げAI","台本テキストを音声に変換して映像に追加","CapCut TTS"],
         ["字幕AI","ナレーション音声から字幕を自動生成","CapCut AI字幕"],
         ["人間","字幕修正・クリップ編集・BGM追加・完成品の書き出し","CapCut 手動編集"]],
        Inches(0.5), Inches(2.0), Inches(9), Inches(3.3),
        col_widths=[0.20, 0.40, 0.40])

    # ── Section 1：イントロ ──
    slide_section(prs, "1", "イントロ", "10分", "CapCutにログインして基本画面を確認")
    sl = blank_slide(prs)
    slide_header(sl, "イントロ — CapCut基本画面の説明", "10分")
    add_table(sl, ["確認項目", "内容"],
        [["CapCutログイン","capcut.com または CapCutアプリ（PC版推奨）でTikTokアカウント or Google / Appleで登録"],
         ["新規プロジェクト","「新規作成」→「ビデオ」を選択して編集画面へ"],
         ["Seedance","「AI」メニュー内 → テキストから映像を生成する機能（クレジット制）"],
         ["TTS（テキスト読み上げ）","「テキスト」→「テキスト読み上げ」で台本を音声に変換"],
         ["AI字幕","「字幕」→「自動字幕生成」で音声から字幕を自動生成（月2回まで無料）"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.6),
        col_widths=[0.28, 0.72])
    info_box(sl, "Seedanceはクレジット制です。本日は1人あたり○回を上限にしてください（講師より案内）",
             Inches(0.5), Inches(4.63), Inches(9), Inches(0.65),
             bg=RGBColor(0xFF, 0xF0, 0xE8), text_color=C_ACCENT, size=12)

    # ── Section 2：AI動画生成 ──
    slide_section(prs, "2", "AI動画生成", "15分", "テキスト指示だけで映像クリップが生成される体験")
    sl = blank_slide(prs)
    slide_header(sl, "AI動画生成 — CapCut Seedance の使い方", "デモ5分 + 実践10分")
    add_table(sl, ["Step", "内容"],
        [["1 Seedanceを開く","CapCut「AI」メニュー → 「Seedance」を選択"],
         ["2 テキストを入力","「どんな映像にしたいか」を日本語で説明 → 生成（30秒〜1分程度かかる）"],
         ["3 映像を確認","生成された映像クリップを再生して確認 → 気に入ったらプロジェクトに追加"],
         ["4 繰り返し","別のシーンを追加したい場合は再度テキストを入力して生成"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.8),
        col_widths=[0.25, 0.75])
    tb(sl, "生成テーマ例（自由に選んでください）：",
       Inches(0.5), Inches(3.85), Inches(9), Inches(0.4),
       size=13, bold=True, color=C_TEXT)
    examples = sl.shapes.add_textbox(Inches(0.5), Inches(4.28), Inches(9), Inches(0.8))
    tf = examples.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "会社・お店の雰囲気紹介 ／ 商品・サービスのイメージ映像 ／ 季節のご挨拶 ／ 採用広報（職場の一日）"
    r.font.size = Pt(12)
    r.font.color.rgb = C_MUTED

    # ── Section 3：ナレーション ──
    slide_section(prs, "3", "ナレーション追加", "15分", "Geminiで台本を作り → CapCut TTSで音声に変換")
    sl = blank_slide(prs)
    slide_header(sl, "ナレーション追加 — Gemini で台本 → CapCut TTS で音声化", "15分")
    add_table(sl, ["Step", "内容", "ツール"],
        [["1 Geminiを開く","gemini.google.com にアクセス（Googleアカウントでログイン）","Gemini"],
         ["2 台本を生成","ひな型でプロンプトを入力 → 生成された台本テキストをコピー","Gemini"],
         ["3 TTSに貼り付け","CapCut「テキスト」→「テキスト読み上げ」→ 台本をペースト","CapCut"],
         ["4 声を選ぶ","日本語の音声種類（男性・女性・落ち着き・明るい等）を選択","CapCut"],
         ["5 映像に追加","生成された音声をタイムラインに追加・長さを調整","CapCut"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.5),
        col_widths=[0.22, 0.55, 0.23])
    info_box(sl, "台本ひな型：「[テーマ]の紹介動画のナレーション原稿を1分で読める長さで作ってください」",
             Inches(0.5), Inches(4.52), Inches(9), Inches(0.65),
             text_color=C_TEXT, size=13)

    # ── Section 4：AI自動字幕 ──
    slide_section(prs, "4", "AI自動字幕", "10分", "ナレーション音声から字幕を自動生成")
    sl = blank_slide(prs)
    slide_header(sl, "AI自動字幕 — 生成して確認まで（修正は5で）", "10分")
    add_table(sl, ["Step", "内容"],
        [["1 字幕メニューを開く","CapCut「字幕」→「自動字幕生成」を選択"],
         ["2 言語を設定","日本語を選択 → 「生成」をクリック（少し時間がかかる）"],
         ["3 結果を確認","ナレーション音声が正しく文字に変換されているか確認"],
         ["4 修正は次のセクションで","ここでは生成結果の確認まで。修正はセクション5（基本編集）で行う"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.9),
        col_widths=[0.28, 0.72])
    info_box(sl, "注意：AI自動字幕は無料版では月2回まで。事前に講師側で確認しています",
             Inches(0.5), Inches(3.95), Inches(9), Inches(0.65),
             bg=RGBColor(0xFF, 0xF0, 0xE8), text_color=C_ACCENT, size=13)
    tb(sl, "キャプションテンプレート（文字装飾）は時間があれば口頭で紹介します",
       Inches(0.5), Inches(4.75), Inches(9), Inches(0.5),
       size=12, italic=True, color=C_MUTED)

    # ── 休憩 ──
    slide_break(prs)

    # ── Section 5：基本編集・書き出し ──
    slide_section(prs, "5", "基本編集・書き出し", "50分", "AIの素材を人間が磨いて完成品にする")
    sl = blank_slide(prs)
    slide_header(sl, "基本編集 — 字幕修正 → クリップ編集 → 書き出し", "50分")
    add_table(sl, ["作業", "内容", "時間"],
        [["字幕修正","AI字幕のテキストを確認 → 誤認識・読み間違いを手動で修正","〜15分"],
         ["クリップ編集","タイミング調整・不要部分カット・テキスト追加・BGM追加など","〜25分"],
         ["書き出し","「書き出し」→ 解像度設定（1080p推奨）→ MP4で保存（1分動画は1〜3分程度）","10分"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.5),
        col_widths=[0.22, 0.55, 0.23])
    tb(sl, "基本編集の意義：AIが作った素材を人間が磨く — これが「完成品への最終仕上げ」",
       Inches(0.5), Inches(3.55), Inches(9), Inches(0.55),
       size=14, bold=True, color=C_PRIMARY)
    info_box(sl, "早く終わった方はセクション2に戻って2本目を作成してください",
             Inches(0.5), Inches(4.28), Inches(9), Inches(0.65),
             bg=C_GREEN_BG, text_color=C_GREEN, size=13)

    # ── Section 6：まとめ ──
    slide_section(prs, "6", "まとめ", "10分", "ダウンロード・SNS投稿・活用アドバイス")
    sl = blank_slide(prs)
    slide_header(sl, "まとめ — 活用アドバイス＆サービス案内", "10分")
    add_table(sl, ["活用ポイント", "内容"],
        [["ダウンロード","スマホ・PCへの保存 → スマホはCapCutアプリからも保存可"],
         ["SNS投稿","CapCutからTikTok・Instagramへ直接投稿できる"],
         ["後から編集","CapCutアカウントにログインすれば続きから編集・複製して再利用"],
         ["量産のコツ","テーマを変えれば今日と同じ流れで次の動画が作れる"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.8),
        col_widths=[0.25, 0.75])
    rect(sl, Inches(0.5), Inches(3.85), Inches(9), Inches(1.55), C_BG_LIGHT)
    tb(sl, "次のステップ | ざつね屋のサービス案内",
       Inches(0.7), Inches(3.95), Inches(8.5), Inches(0.45),
       size=13, bold=True, color=C_PRIMARY)
    svc = sl.shapes.add_textbox(Inches(0.7), Inches(4.42), Inches(8.5), Inches(0.85))
    tf = svc.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "📚 実践編（画像系）でCanvaのデザイン力をさらに高める"
    r.font.size = Pt(12)
    r.font.color.rgb = C_TEXT
    add_line(tf, "🛠 AI業務改善オーダーメイド・🤝 AI開発伴走 — 自社の課題を一緒に解決する", size=12, color=C_TEXT)

    prs.save(output_path)
    print(f"[OK] course04: {output_path} ({len(prs.slides)} slides)")


# ═══════════════════════════════════════════════════════════
#  コース⑤：実践編（Claude Code特化）
# ═══════════════════════════════════════════════════════════

def make_course05(output_path):
    prs = new_prs()

    # ── タイトル ──
    sl = blank_slide(prs)
    set_bg(sl, C_PRIMARY)
    rect(sl, Inches(0), Inches(3.7), SLIDE_W, Inches(0.06), C_ACCENT)
    tb(sl, "ざつね屋", Inches(0.6), Inches(0.4), Inches(4), Inches(0.45),
       size=13, color=C_PALE_BLUE)
    tb(sl, "コース ⑤", Inches(0.6), Inches(0.95), Inches(9), Inches(0.55),
       size=18, color=C_LIGHT_BLUE)
    tb(sl, "実践編（Claude Code特化）", Inches(0.6), Inches(1.5), Inches(9), Inches(1.0),
       size=38, bold=True, color=C_WHITE)
    tb(sl, "コード経験不要。社内の困りごとをAIで解決する仕組みを自分で作る 120分",
       Inches(0.6), Inches(2.6), Inches(9), Inches(0.7),
       size=15, color=C_LIGHT_BLUE)
    tb(sl, "ツール：Claude Code デスクトップ版　｜　定員：個別1名 / グループMAX3名",
       Inches(0.6), Inches(4.1), Inches(9), Inches(0.5),
       size=13, color=C_PALE_BLUE)

    # ── アジェンダ ──
    sl = blank_slide(prs)
    slide_header(sl, "本日のアジェンダ")
    add_table(sl, ["#", "セクション", "内容", "時間"],
        [["1","Claude Codeとは","4軸マップ・他AIとの違い・できること/できないこと","20分"],
         ["2","セットアップ確認","デスクトップ版起動・サインイン・動作確認","30分"],
         ["3","ユースケース","フォルダ整理・CSV整形・定型文書生成（各3分）","10分"],
         ["—","休憩","","10分"],
         ["4","応用ワーク","CLAUDE.md作成・テンプレート化・自社業務試行","40分"],
         ["5","まとめ・Q&A","伴走サービス紹介・質疑応答","10分"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.55),
        col_widths=[0.06, 0.22, 0.52, 0.20])

    # ── Section 1：Claude Codeとは ──
    slide_section(prs, "1", "Claude Code とは", "20分",
                  "4軸マップでの位置確認 ＆ 他ツールとの違い")
    sl = blank_slide(prs)
    slide_header(sl, "4軸AIマップ — Claude Codeの位置", "3分")
    add_table(sl, ["領域", "対象", "主な用途", "代表ツール"],
        [["文書・テキスト", "全スタッフ", "メール・報告書・広報文", "ChatGPT / Gemini / Claude"],
         ["画像・ビジュアル", "広報・SNS担当", "SNS画像・チラシ", "Canva / Firefly"],
         ["動画・映像", "広報・採用担当", "SNS動画・採用動画", "CapCut / Canva"],
         ["自動化・開発 ← ここ", "経営者・業務改善担当", "繰り返し作業の自動化・仕組み化", "Claude Code"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.85),
        col_widths=[0.25, 0.22, 0.30, 0.23])
    info_box(sl, "「自動化・開発」領域だけがAPIキーや環境設定が必要。Claude Codeはその敷居を最も下げたツールです",
             Inches(0.5), Inches(3.88), Inches(9), Inches(0.65),
             text_color=C_TEXT, size=13)

    sl = blank_slide(prs)
    slide_header(sl, "Claude と Claude Code の違い", "4分")
    add_table(sl, ["", "Claude（claude.ai）", "Claude Code"],
        [["種類", "会話型AI", "AIエージェント"],
         ["使い方", "ブラウザでチャット", "デスクトップアプリ"],
         ["できること", "答える・書く・提案する", "ファイルを作る・動かす・整理する"],
         ["イメージ", "優秀な相談相手", "実際に手を動かす部下"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.15),
        col_widths=[0.18, 0.41, 0.41])
    info_box(sl, "Claude Codeは「答えを出す」だけでなく「実際にファイルを操作・生成する」ところが決定的な違いです",
             Inches(0.5), Inches(4.18), Inches(9), Inches(0.65),
             text_color=C_TEXT, size=13)

    sl = blank_slide(prs)
    slide_header(sl, "Claude Code — できること / できないこと", "4分")
    add_table(sl, ["できること", "できないこと"],
        [["ファイル・フォルダ操作", "ネット上の操作（注文・送信・フォーム入力など）"],
         ["CSV・データの整形・集計", "完全な自動実行（毎回指示は必要）"],
         ["定型文書の一括生成", "複雑なシステム連携（別途開発が必要）"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.3),
        col_widths=[0.42, 0.58])
    tb(sl, "利用条件",
       Inches(0.5), Inches(3.3), Inches(9), Inches(0.4),
       size=14, bold=True, color=C_TEXT)
    add_table(sl, ["必要なもの", "詳細"],
        [["PC（Windows / Mac）", "タブレット・スマホは不可"],
         ["Claude Code デスクトップ版", "無料インストール（claude.ai/download）"],
         ["Anthropicアカウント + Proプラン", "$20/月・約3,000円・APIキー不要（対話的利用のため）"]],
        Inches(0.5), Inches(3.78), Inches(9), Inches(1.6),
        col_widths=[0.35, 0.65])

    sl = blank_slide(prs)
    slide_header(sl, "APIとは — 「そこまでやりたい方は伴走サービスへ」の導線", "10分")
    add_table(sl, ["用途", "必要なもの", "例"],
        [["対話的利用（本日の研修）","Anthropicアカウント + Proプラン","ファイル整理・CSV整形・文書生成"],
         ["開発・システム連携","APIキー（別途従量課金）","Excelボタンで動く仕組み・LINE Bot・自社アプリ組み込み"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.2),
        col_widths=[0.28, 0.32, 0.40])
    tb(sl, "ハルシネーション注意",
       Inches(0.5), Inches(3.2), Inches(9), Inches(0.4),
       size=13, bold=True, color=C_ACCENT)
    info_box(sl, "Claude Codeも自信満々に間違えます。重要なファイル操作の前は必ず「何をするか教えてから動いてください」と指示を加えましょう。CLAUDE.mdで防止策を書いておくのが有効です（セクション4で実践します）",
             Inches(0.5), Inches(3.68), Inches(9), Inches(1.2),
             bg=RGBColor(0xFF, 0xF0, 0xE8), text_color=C_TEXT, size=12)

    # ── Section 2：セットアップ確認 ──
    slide_section(prs, "2", "セットアップ確認", "30分", "デスクトップ版起動・サインイン・動作確認")
    sl = blank_slide(prs)
    slide_header(sl, "セットアップ確認 — 全員動作確認まで", "30分")
    add_table(sl, ["Step", "内容", "詰まったら"],
        [["1 起動","Claude Codeデスクトップ版を起動（未インストールの場合はその場でインストール）","講師が個別サポート"],
         ["2 サインイン","AnthropicアカウントでGoogle または メールアドレスでログイン","アカウント未作成の場合は作成"],
         ["3 Proプラン確認","設定画面またはプラン情報でProプラン以上が有効か確認","未加入の場合は当日申し込み"],
         ["4 動作テスト","日本語で「こんにちは、今日はよろしくお願いします」と話しかけて返答が来るまで確認","アプリ再起動 → 再サインイン"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.1),
        col_widths=[0.20, 0.52, 0.28])
    tb(sl, "3名以下なので詰まった方を個別サポートします。遠慮なく声をかけてください",
       Inches(0.5), Inches(5.12), Inches(9), Inches(0.35),
       size=12, italic=True, color=C_MUTED)

    # ── Section 3：ユースケース ──
    slide_section(prs, "3", "ユースケース", "10分", "3つの実用シーンを各3分でデモ体験")
    sl = blank_slide(prs)
    slide_header(sl, "ユースケース ①：フォルダ整理", "約3分")
    info_box(sl, "指示例：「Downloadsにある画像・PDF・Wordを種類ごとにフォルダ分けして」",
             Inches(0.5), Inches(0.85), Inches(9), Inches(0.65),
             text_color=C_TEXT, size=14, bold=True)
    tb(sl, "最初にやる理由：持参PCの自分のフォルダが整理される体験が最もインパクト大でシンプル",
       Inches(0.5), Inches(1.65), Inches(9), Inches(0.55),
       size=13, color=C_MUTED, italic=True)
    tb(sl, "Claude Codeが実際にファイルを動かしてくれる様子を確認してください",
       Inches(0.5), Inches(2.3), Inches(9), Inches(0.55),
       size=14, color=C_TEXT)
    info_box(sl, "確認ポイント：実行前に「どのフォルダを作りますか？」と確認してから動くかどうか",
             Inches(0.5), Inches(3.05), Inches(9), Inches(0.65),
             bg=C_BG_LIGHT, text_color=C_PRIMARY, size=13)

    sl = blank_slide(prs)
    slide_header(sl, "ユースケース ②③：CSV整形 ＆ 定型文書生成", "各約3分")
    add_table(sl, ["ユースケース", "指示例"],
        [["② CSV・データ整形","「顧客リストから氏名・メールを抜き出して50音順に並べ直して」"],
         ["③ 定型文書の一括生成","「このリストの各社名に向けた個別のお礼メールをテキストファイルで作って」"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.1),
        col_widths=[0.28, 0.72])
    tb(sl, "3つのユースケースで共通していること",
       Inches(0.5), Inches(3.15), Inches(9), Inches(0.45),
       size=14, bold=True, color=C_PRIMARY)
    info_box(sl, "「繰り返し作業」を自動化している。一度指示を覚えれば次回は数秒で完了する",
             Inches(0.5), Inches(3.68), Inches(9), Inches(0.65),
             text_color=C_TEXT, size=13, bold=True)

    # ── 休憩 ──
    slide_break(prs)

    # ── Section 4：応用ワーク ──
    slide_section(prs, "4", "応用ワーク", "40分", "CLAUDE.md作成 → テンプレート化 → 自社業務試行")
    sl = blank_slide(prs)
    slide_header(sl, "CLAUDE.md とは — 作成実践", "15分")
    add_table(sl, ["カテゴリ", "記載例"],
        [["事業者情報","会社名・業種・拠点"],
         ["出力ルール","「出力は必ず日本語で」「箇条書き優先」"],
         ["保存先","「ファイルはDesktopに保存すること」"],
         ["禁止事項","「個人情報・機密情報は扱わない」"],
         ["よく使う形式","「文書はWord(.docx)・データはCSVで出力」"],
         ["口調","「敬語・ですます調で出力」"],
         ["ハルシネーション対策","「わからない場合は作業を止めて確認すること」「実行前に何をするか説明してから動くこと」"]],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.3),
        col_widths=[0.28, 0.72])

    sl = blank_slide(prs)
    slide_header(sl, "テンプレート化 ＆ AIが止まるケースと対処", "5分 + 5分")
    tb(sl, "テンプレート化とは",
       Inches(0.5), Inches(0.82), Inches(4.5), Inches(0.4),
       size=13, bold=True, color=C_PRIMARY)
    temp_box = sl.shapes.add_textbox(Inches(0.5), Inches(1.25), Inches(4.5), Inches(1.2))
    tf = temp_box.text_frame
    tf.word_wrap = True
    for i, t in enumerate(["よく使う指示文を記録しておく機能",
                            "同じ作業を対象を変えて使い回せる",
                            "CLAUDE.mdに入れておけば名前で呼び出せる"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"・{t}"
        r.font.size = Pt(12)
        r.font.color.rgb = C_TEXT
    add_table(sl, ["トラブル", "対処"],
        [["返答が途中で止まった","「続けてください」と入力"],
         ["意図と違う処理をした","「待って」と入力 → 意図を正確に伝え直す"],
         ["エラーメッセージが出た","メッセージをそのまま貼り付けて「これはどういう意味ですか？」と聞く"],
         ["コンテキストが長くなった","「ここまでの作業内容をまとめてください」→ 新規セッションに貼り付け"],
         ["応答しない・フリーズ","数分待ってアプリを再起動して再試行"]],
        Inches(5.0), Inches(0.82), Inches(4.8), Inches(4.55),
        col_widths=[0.40, 0.60])

    sl = blank_slide(prs)
    slide_header(sl, "自社業務への試行", "10分")
    tb(sl, "効率化したい作業を Claude Code に相談しながら自由に試してください",
       Inches(0.5), Inches(0.88), Inches(9), Inches(0.55),
       size=14, color=C_TEXT)
    add_table(sl, ["困ったら", "対処法"],
        [["テーマが思いつかない","セクション3のユースケースをさらに深める（自社のデータで試す）"],
         ["エラーが出た","エラーメッセージを貼り付けて「これはどういう意味ですか？」"],
         ["完成しなくて大丈夫","「設計まで到達できれば十分」— 完成が目標ではない"],
         ["どうしてもわからない","ブラウザ版Claude（claude.ai）に状況を説明して解決策を聞く"]],
        Inches(0.5), Inches(1.6), Inches(9), Inches(2.8),
        col_widths=[0.30, 0.70])
    info_box(sl, "時間いっぱいAIと対話してください。時間が来たら作業途中でOKです",
             Inches(0.5), Inches(4.6), Inches(9), Inches(0.65),
             bg=C_GREEN_BG, text_color=C_GREEN, size=13)

    # ── Section 5：まとめ・Q&A ──
    slide_section(prs, "5", "まとめ・Q&A", "10分", "伴走サービス紹介 ＆ 質疑応答")
    sl = blank_slide(prs)
    slide_header(sl, "まとめ — 本日の学び＆次のステップ", "10分")
    learns = sl.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(5.0), Inches(2.2))
    tf = learns.text_frame
    tf.word_wrap = True
    for i, item in enumerate([
        "[OK] Claude Codeが「実際に手を動かす」AIだとわかった",
        "[OK] フォルダ整理・CSV整形・文書生成を体験した",
        "[OK] CLAUDE.mdで指示を記録・再利用できるとわかった",
        "[OK] トラブル時の対処法を持ち帰れた",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = item
        r.font.size = Pt(13)
        r.font.color.rgb = C_TEXT

    rect(sl, Inches(5.3), Inches(0.85), Inches(4.5), Inches(2.2), C_BG_LIGHT)
    tb(sl, "AI開発伴走サービス",
       Inches(5.5), Inches(0.95), Inches(4.1), Inches(0.45),
       size=13, bold=True, color=C_PRIMARY)
    banso = sl.shapes.add_textbox(Inches(5.5), Inches(1.45), Inches(4.1), Inches(1.5))
    tf2 = banso.text_frame
    tf2.word_wrap = True
    for i, item in enumerate([
        "「もっと深めたい」「自社業務に組み込みたい」",
        "ざつね屋と共同開発 → MTGで知見を統合",
        "月1回プラン：60,000円/月",
        "月2回プラン：100,000円/月",
    ]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        r = p.add_run()
        r.text = item
        r.font.size = Pt(12)
        r.font.color.rgb = C_TEXT

    tb(sl, "Q&A — 質問・相談はなんでもどうぞ",
       Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
       size=14, bold=True, color=C_TEXT)
    tb(sl, "「わからないことはAIに聞いてみてください」",
       Inches(0.5), Inches(3.85), Inches(9), Inches(0.6),
       size=15, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    tb(sl, "本日はありがとうございました",
       Inches(0.5), Inches(4.55), Inches(9), Inches(0.5),
       size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    print(f"[OK] course05: {output_path} ({len(prs.slides)} slides)")


# ─────────────────────────────────────────────────────────
if __name__ == "__main__":
    make_course02("training-materials/course02_jissen_bunsho.pptx")
    make_course03("training-materials/course03_jissen_gazo.pptx")
    make_course04("training-materials/course04_jissen_doga.pptx")
    make_course05("training-materials/course05_jissen_claude_code.pptx")
    print("[DONE] 4 courses generated.")
