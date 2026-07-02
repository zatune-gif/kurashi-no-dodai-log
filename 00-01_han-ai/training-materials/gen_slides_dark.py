"""
研修スライド PDF 生成スクリプト v4
"""
import sys, re, asyncio, os, html as htmlmod
from pptx import Presentation
from playwright.async_api import async_playwright

COURSES = [
    ("course01_AI活用知識編.pptx",        "course01_AI活用知識編.pdf"),
    ("course02_jissen_bunsho.pptx",       "course02_jissen_bunsho.pdf"),
    ("course03_jissen_gazo.pptx",         "course03_jissen_gazo.pdf"),
    ("course04_jissen_doga.pptx",         "course04_jissen_doga.pdf"),
    ("course05_jissen_claude_code.pptx",  "course05_jissen_claude_code.pdf"),
]

# ── デザイントークン ──────────────────────────────────────────
BG       = "rgb(28,28,28)"
BG_CARD  = "rgb(37,37,37)"
BG_PAGE  = "rgb(17,17,17)"
ACCENT   = "rgb(218,119,86)"
DIVIDER  = "rgb(56,56,56)"
TEXT     = "rgb(238,238,238)"
MUTED    = "rgb(144,144,144)"
DARK     = "rgb(28,28,28)"

# ── ブロック管理 ──────────────────────────────────────────────
BLOCK_HEADERS = {
    "社内ルール（最低限）", "出力を使うときの注意", "本日の学び",
    "次のステップ | ざつね屋のサービス案内", "実践（30分） — ツール",
}
SKIP_BLOCKS = {"持ち帰り資料"}  # このブロックとその内容は非表示

# ── テキストユーティリティ ────────────────────────────────────
_EMOJI_RE = re.compile(
    r'[\U0001F000-\U0001FFFF'   # 絵文字全般
    r'☀-➿'            # 記号・装飾
    r'⌀-⏿'            # 技術記号
    r'⬀-⯿'            # 矢印・図形
    r'─-◿'            # ボックス描画
    r']+',
    flags=re.UNICODE,
)

def strip_emoji(s):
    """絵文字・アイコン文字を除去"""
    return _EMOJI_RE.sub('', str(s)).strip()

def is_subitem(s):
    """→ で始まるサブ項目かどうか"""
    return s.startswith('→') or s.startswith('⇒')

def clean(s):
    """絵文字除去 → strip"""
    return strip_emoji(s)

def h(s):
    return htmlmod.escape(str(s))

def filter_excluded(items):
    """SKIP_BLOCKSブロックとその内容を除去"""
    result = []
    skipping = False
    for t in items:
        if t in SKIP_BLOCKS:
            skipping = True
            continue
        if skipping and t in BLOCK_HEADERS:
            skipping = False
        if not skipping:
            result.append(t)
    return result

# ── ベースCSS ────────────────────────────────────────────────
BASE_CSS = """
@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+JP:wght@400;500;700;800;900&display=swap');
*, *::before, *::after { box-sizing: border-box; -webkit-print-color-adjust: exact; print-color-adjust: exact; }
html, body { margin: 0; padding: 0; background: rgb(17,17,17); }
body { font-family: 'Noto Sans JP', 'Yu Gothic', 'Hiragino Sans', sans-serif; }
h1, h2, h3, p { margin: 0; padding: 0; }
@page { size: 1280px 720px; margin: 0; }
@media print {
  html, body { background: #ffffff; }
  .deck-surround { padding: 0 !important; gap: 0 !important; background: #ffffff !important; }
}
"""

def footer(top=False):
    pos = "top:28px;right:80px;" if top else "bottom:56px;right:80px;"
    return (
        f'<div style="position:absolute;{pos}display:flex;align-items:center;'
        f'gap:12px;font-size:20px;font-weight:700;color:{TEXT};">'
        f'<span style="width:10px;height:10px;border-radius:50%;background:{ACCENT};'
        f'display:inline-block;"></span>ざつね屋</div>'
    )

def top_bar():
    return f'<div style="position:absolute;top:0px;left:0px;width:100%;height:8px;background:{ACCENT};"></div>'

# ── テキスト抽出 ────────────────────────────────────────────
def get_texts(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    out.append(t)
    return out

def get_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None

# ── スライドタイプ判定 ──────────────────────────────────────
def detect_type(slide):
    texts = get_texts(slide)
    if not texts:
        return "blank"
    arrows = sum(1 for t in texts if t in ("↑","↓","←","→"))
    if arrows >= 3:
        return "grid"
    if texts[0] == "ざつね屋":
        return "cover"
    if any("Before" in t or "After" in t for t in texts):
        return "before_after"
    if re.match(r'^[1-9]$', texts[0]):
        return "section"
    if get_table(slide):
        return "table"
    return "bullet"

# ── 表紙 ────────────────────────────────────────────────────
def render_cover(slide):
    texts = get_texts(slide)
    badge    = texts[1] if len(texts) > 1 else ""
    title    = texts[2] if len(texts) > 2 else ""
    subtitle = texts[3] if len(texts) > 3 else ""
    return (
        f'<section style="width:1280px;height:720px;background:{BG};color:{TEXT};'
        f'box-sizing:border-box;position:relative;overflow:hidden;break-after:page;'
        f'padding:80px;display:flex;flex-direction:column;justify-content:center;">'
        f'{top_bar()}'
        f'<div style="display:inline-flex;align-self:flex-start;align-items:center;'
        f'background:{ACCENT};color:{DARK};padding:10px 24px;border-radius:6px;'
        f'font-weight:800;font-size:22px;letter-spacing:0.04em;">{h(badge)}</div>'
        f'<h1 style="margin:36px 0px 0px;font-size:92px;font-weight:900;'
        f'line-height:1.08;letter-spacing:0.01em;">{h(title)}</h1>'
        f'<p style="margin:32px 0px 0px;font-size:28px;font-weight:500;'
        f'color:{MUTED};line-height:1.5;">{h(subtitle)}</p>'
        f'{footer(top=True)}</section>'
    )

# ── セクション区切り ─────────────────────────────────────────
def render_section(slide):
    texts = get_texts(slide)
    num   = texts[0] if len(texts) > 0 else ""
    name  = texts[1] if len(texts) > 1 else ""
    # texts[2] は時間（アジェンダ以外では非表示）
    desc  = texts[3] if len(texts) > 3 else ""
    return (
        f'<section style="width:1280px;height:720px;background:{BG};color:{TEXT};'
        f'box-sizing:border-box;position:relative;overflow:hidden;break-after:page;'
        f'padding:80px;display:flex;align-items:center;gap:72px;">'
        f'<div style="font-size:340px;font-weight:900;line-height:0.8;color:{ACCENT};flex-shrink:0;">{h(num)}</div>'
        f'<div style="width:3px;align-self:stretch;margin:64px 0px;background:{DIVIDER};flex-shrink:0;"></div>'
        f'<div style="display:flex;flex-direction:column;gap:28px;">'
        f'<h2 style="margin:0px;font-size:60px;font-weight:800;line-height:1.15;">{h(name)}</h2>'
        f'<p style="margin:0px;font-size:26px;font-weight:500;color:{MUTED};'
        f'line-height:1.6;max-width:640px;">{h(desc)}</p>'
        f'</div>{footer(top=True)}</section>'
    )

# ── コンテンツラッパー ──────────────────────────────────────
def content_wrap(title_text, inner_html):
    return (
        f'<section style="width:1280px;height:720px;background:{BG};color:{TEXT};'
        f'box-sizing:border-box;position:relative;overflow:hidden;break-after:page;'
        f'padding:64px 80px;display:flex;flex-direction:column;">'
        f'{footer(top=True)}'
        f'<h2 style="margin:0px 0px 28px;font-size:48px;font-weight:800;'
        f'padding-bottom:22px;border-bottom:4px solid {ACCENT};">{h(title_text)}</h2>'
        f'{inner_html}</section>'
    )

# ── 箇条書き ─────────────────────────────────────────────────
def render_bullet(slide):
    texts = get_texts(slide)
    if not texts:
        return f'<section style="break-after:page;width:1280px;height:720px;background:{BG};"></section>'

    title = texts[0]
    items_all = filter_excluded(texts[1:])

    if any(t in BLOCK_HEADERS for t in items_all):
        return render_two_col(slide)

    # 注記テキストを分離（先頭から見てあきらかに補足のもの）
    note_pfx = ("⚠","💡","❗","📋","📚","🛠","✅","注意","ポイント","次のステップ")
    mains = [t for t in items_all if not any(t.startswith(p) for p in note_pfx)]
    notes = [t for t in items_all if any(t.startswith(p) for p in note_pfx)]

    n = len(mains)
    sparse = n <= 2
    fs = "32px" if n <= 5 else "24px" if n <= 7 else "20px"
    container_extra = "justify-content:center;gap:32px;" if sparse else ""

    rows = ""
    for i, item in enumerate(mains):
        c = clean(item)
        if is_subitem(c):
            text = re.sub(r'^[→⇒]+\s*', '', c)
            border = ""
            rows += (
                f'<div style="display:flex;align-items:flex-start;gap:20px;'
                f'padding:6px 0 6px 44px;">'
                f'<span style="font-size:18px;font-weight:500;color:{MUTED};'
                f'line-height:1.6;">{h(text)}</span></div>'
            )
        else:
            border = (f"border-bottom:1px solid {DIVIDER};" if i < n - 1 else "") if not sparse else ""
            pad = "padding:20px 0;" if sparse else ""
            flex = "display:flex;" if sparse else "flex:1 1 0%;display:flex;"
            rows += (
                f'<div style="{flex}align-items:center;gap:28px;{border}{pad}">'
                f'<span style="width:16px;height:16px;background:{ACCENT};flex-shrink:0;"></span>'
                f'<span style="font-size:{fs};font-weight:500;">{h(c)}</span></div>'
            )

    note_html = ""
    if notes:
        note_text = "　".join(clean(n) for n in notes[:2])
        note_html = (
            f'<div style="margin-top:16px;padding:12px 20px;border-left:4px solid {ACCENT};'
            f'background:{BG_CARD};font-size:18px;color:{MUTED};line-height:1.5;">'
            f'{h(note_text)}</div>'
        )

    inner = (
        f'<div style="flex:1 1 0%;display:flex;flex-direction:column;margin-top:8px;{container_extra}">'
        f'{rows}</div>{note_html}'
    )
    return content_wrap(title, inner)

# ── 2カラム ──────────────────────────────────────────────────
def render_two_col(slide):
    texts = get_texts(slide)
    if not texts:
        return f'<section style="break-after:page;width:1280px;height:720px;background:{BG};"></section>'

    title = texts[0]
    rest  = filter_excluded(texts[1:])

    blocks = []
    cur_hd, cur_items = None, []
    for t in rest:
        if t in BLOCK_HEADERS:
            if cur_hd is not None or cur_items:
                blocks.append((cur_hd, cur_items))
            cur_hd, cur_items = t, []
        else:
            cur_items.append(t)
    blocks.append((cur_hd, cur_items))

    cols = ""
    for bh, bi in blocks[:4]:
        items_html = ""
        for it in bi:
            c = clean(it)
            if is_subitem(c):
                text = re.sub(r'^[→⇒]+\s*', '', c)
                items_html += (
                    f'<div style="display:flex;align-items:flex-start;gap:10px;'
                    f'padding-left:20px;font-size:17px;font-weight:400;'
                    f'line-height:1.5;color:{MUTED};">'
                    f'<span style="flex-shrink:0;margin-top:4px;color:{ACCENT};font-size:14px;">╴</span>'
                    f'<span>{h(text)}</span></div>'
                )
            else:
                items_html += (
                    f'<div style="display:flex;align-items:flex-start;gap:12px;'
                    f'font-size:19px;font-weight:500;line-height:1.5;">'
                    f'<span style="width:10px;height:10px;background:{ACCENT};'
                    f'flex-shrink:0;margin-top:5px;"></span>'
                    f'<span>{h(c)}</span></div>'
                )
        cols += (
            f'<div style="background:{BG_CARD};border-radius:8px;padding:22px 26px;'
            f'display:flex;flex-direction:column;gap:10px;">'
            f'<div style="font-size:19px;font-weight:800;color:{ACCENT};'
            f'padding-bottom:10px;border-bottom:1px solid {DIVIDER};">{h(bh or "")}</div>'
            f'{items_html}</div>'
        )

    inner = (
        f'<div style="flex:1 1 0%;display:grid;grid-template-columns:1fr 1fr;gap:20px;">'
        f'{cols}</div>'
    )
    return content_wrap(title, inner)

# ── テーブル ─────────────────────────────────────────────────
def render_table(slide):
    texts = get_texts(slide)
    tbl   = get_table(slide)
    title = texts[0] if texts else ""

    if tbl is None:
        return render_bullet(slide)

    n_rows = len(tbl.rows)
    n_cols = len(tbl.columns)

    compact = n_rows > 5 or n_cols > 3
    font_h  = "22px" if compact else "30px"
    font_b  = "18px" if compact else "26px"
    pad     = "0 20px" if compact else "0 40px"

    # 列0の幅：long row headers に対応するため auto/min-content を使用
    col0_w = "minmax(160px,max-content)" if compact else "minmax(200px,max-content)"
    if n_cols == 2:
        grid_cols = f"1fr 2fr"
    elif n_cols == 3:
        grid_cols = f"{col0_w} 1fr 1fr"
    elif n_cols == 4:
        grid_cols = f"{col0_w} 1fr 1fr 1fr"
    else:
        grid_cols = f"{col0_w} " + " ".join(["1fr"] * (n_cols - 1))

    grid_rows = "auto " + " ".join(["1fr"] * max(n_rows - 1, 1))

    cells = ""
    for r in range(n_rows):
        for c in range(n_cols):
            ct = h(clean(tbl.cell(r, c).text.strip()))
            if r == 0:
                if c == 0:
                    ct = ""
                    st = f"background:{BG};display:flex;align-items:center;justify-content:center;"
                else:
                    st = (
                        f"background:{ACCENT};color:{DARK};display:flex;"
                        f"align-items:center;justify-content:center;"
                        f"font-size:{font_h};font-weight:800;text-align:center;"
                    )
            elif c == 0:
                st = (
                    f"background:{BG_CARD};display:flex;align-items:center;"
                    f"justify-content:center;padding:{pad};font-size:{font_b};"
                    f"font-weight:800;color:{ACCENT};text-align:center;"
                )
            else:
                st = (
                    f"background:{BG_CARD};display:flex;align-items:center;"
                    f"padding:{pad};font-size:{font_b};font-weight:500;line-height:1.5;"
                )
            cells += f'<div style="{st}">{ct}</div>'

    tbl_texts = set()
    for r in range(n_rows):
        for c in range(n_cols):
            tbl_texts.add(tbl.cell(r, c).text.strip())
    extra = [t for t in texts[1:] if t not in tbl_texts and len(t) > 10 and "持ち帰り資料" not in t]
    note_html = ""
    if extra:
        note_html = (
            f'<div style="margin-top:10px;padding:10px 18px;border-left:4px solid {ACCENT};'
            f'background:{BG_CARD};font-size:16px;color:{MUTED};line-height:1.5;">'
            f'{h(clean(extra[0]))}</div>'
        )

    inner = (
        f'<div style="flex:1 1 0%;display:grid;grid-template-columns:{grid_cols};'
        f'grid-template-rows:{grid_rows};gap:2px;background:{DIVIDER};border:2px solid {DIVIDER};">'
        f'{cells}</div>{note_html}'
    )
    return content_wrap(title, inner)

# ── 2×2グリッド ─────────────────────────────────────────────
def render_grid(slide):
    texts = get_texts(slide)
    filtered = [t for t in texts if t not in ("↑","↓","←","→") and t.strip()]
    title = filtered[0] if filtered else ""
    cells = filtered[1:]

    nums = ["01","02","03","04"]
    cards = ""
    for i, cell in enumerate(cells[:4]):
        parts   = cell.split("\n", 1)
        c_title = clean(parts[0].strip("【】●").strip())
        c_desc  = clean(parts[1].strip()) if len(parts) > 1 else ""
        num     = nums[i] if i < len(nums) else f"0{i+1}"
        cards += (
            f'<div style="background:{BG_CARD};border-radius:8px;padding:36px 40px;'
            f'display:flex;flex-direction:column;justify-content:center;gap:14px;">'
            f'<span style="font-size:26px;font-weight:900;color:{ACCENT};">{num}</span>'
            f'<h3 style="margin:0px;font-size:38px;font-weight:800;">{h(c_title)}</h3>'
            f'<p style="margin:0px;font-size:22px;font-weight:500;color:{MUTED};line-height:1.5;">{h(c_desc)}</p>'
            f'</div>'
        )

    inner = (
        f'<div style="flex:1 1 0%;display:grid;grid-template-columns:1fr 1fr;'
        f'grid-template-rows:1fr 1fr;gap:24px;">{cards}</div>'
    )
    return content_wrap(title, inner)

# ── Before/After ─────────────────────────────────────────────
def render_before_after(slide):
    texts = get_texts(slide)
    title = texts[0] if texts else ""
    bad_label = good_label = bad_text = good_text = ""
    state = "none"
    for t in texts[1:]:
        if "Before" in t or "❌" in t or "悪い" in t:
            bad_label = clean(t); state = "bad"
        elif "After" in t or "✅" in t or "良い" in t:
            good_label = clean(t); state = "good"
        elif re.match(r'^[↓⬇↑⬆→⇒➡]\s*', t):
            pass  # 矢印・遷移注記はスキップ
        elif state == "bad":
            bad_text += clean(t) + "\n"
        elif state == "good":
            good_text += clean(t) + "\n"

    inner = (
        f'<div style="flex:1 1 0%;display:grid;grid-template-columns:1fr 80px 1fr;gap:0;align-items:center;">'
        f'<div style="background:{BG_CARD};border-radius:8px;padding:28px 32px;'
        f'height:100%;display:flex;flex-direction:column;gap:14px;">'
        f'<div style="font-size:18px;font-weight:700;color:#EF4444;">{h(bad_label)}</div>'
        f'<div style="font-size:22px;font-weight:500;line-height:1.6;white-space:pre-wrap;">{h(bad_text.strip())}</div></div>'
        f'<div style="font-size:32px;color:{ACCENT};text-align:center;">→</div>'
        f'<div style="background:{BG_CARD};border-radius:8px;padding:28px 32px;'
        f'height:100%;display:flex;flex-direction:column;gap:14px;">'
        f'<div style="font-size:18px;font-weight:700;color:#22C55E;">{h(good_label)}</div>'
        f'<div style="font-size:22px;font-weight:500;line-height:1.6;white-space:pre-wrap;">{h(good_text.strip())}</div></div>'
        f'</div>'
    )
    return content_wrap(title, inner)

# ── ルーター ─────────────────────────────────────────────────
def render_slide(slide):
    t = detect_type(slide)
    if t == "cover":        return render_cover(slide)
    if t == "section":      return render_section(slide)
    if t == "table":        return render_table(slide)
    if t == "grid":         return render_grid(slide)
    if t == "before_after": return render_before_after(slide)
    return render_bullet(slide)

def should_skip_slide(slide):
    """タイトル（texts[0]）に除外キーワードを含むスライドのみスキップ"""
    texts = get_texts(slide)
    if not texts:
        return False
    return "持ち帰り資料" in texts[0]

def build_html(pptx_path):
    prs = Presentation(pptx_path)
    slides_html = "\n".join(
        render_slide(s) for s in prs.slides if not should_skip_slide(s)
    )
    return (
        f'<!DOCTYPE html>\n<html><head>\n<meta charset="utf-8">\n'
        f'<style>{BASE_CSS}</style>\n</head>\n<body>\n'
        f'<div class="deck-surround" style="display:flex;flex-direction:column;'
        f'align-items:center;gap:40px;padding:40px;background:{BG_PAGE};">\n'
        f'{slides_html}\n</div>\n</body></html>'
    )

async def html_to_pdf(html_content, pdf_path):
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page()
        await page.set_content(html_content, wait_until="networkidle")
        await page.wait_for_timeout(2000)
        await page.pdf(path=pdf_path, width="1280px", height="720px", print_background=True)
        await browser.close()

async def main():
    base = os.path.dirname(os.path.abspath(__file__))
    for pptx_name, pdf_name in COURSES:
        pptx_path = os.path.join(base, pptx_name)
        pdf_path  = os.path.join(base, pdf_name)
        if not os.path.exists(pptx_path):
            print(f"[SKIP] {pptx_name} not found")
            continue
        print(f"[GEN] {pptx_name} ...", end=" ", flush=True)
        await html_to_pdf(build_html(pptx_path), pdf_path)
        print("done")
    print("\n✅ 全コース PDF 生成完了")

if __name__ == "__main__":
    asyncio.run(main())
