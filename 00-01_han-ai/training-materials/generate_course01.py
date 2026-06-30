"""
コース①：AI活用知識編 スライド生成スクリプト
Google Slides にインポートして使用する .pptx を生成します。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── カラーパレット（ざつね屋ブランドカラー）──────────────
C_PRIMARY    = RGBColor(0x1E, 0x40, 0xAF)
C_ACCENT     = RGBColor(0xDA, 0x77, 0x56)
C_TEXT       = RGBColor(0x1E, 0x29, 0x3B)
C_MUTED      = RGBColor(0x64, 0x74, 0x8B)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_LIGHT   = RGBColor(0xEF, 0xF6, 0xFF)
C_TABLE_ALT  = RGBColor(0xF8, 0xFA, 0xFC)
C_GREEN      = RGBColor(0x05, 0x96, 0x69)
C_RED        = RGBColor(0xDC, 0x26, 0x26)
C_BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
C_LIGHT_BLUE = RGBColor(0xBF, 0xDB, 0xFE)
C_PALE_BLUE  = RGBColor(0x93, 0xC5, 0xFD)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ── 共通ヘルパー ─────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
    else:
        shp.line.fill.background()
    return shp

def text_box(slide, text, left, top, width, height,
             size=14, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb

def add_para(tf, text, size=13, bold=False, italic=False, color=None, spacing_before=None):
    p = tf.add_paragraph()
    if spacing_before:
        p.space_before = spacing_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color

def add_table(slide, headers, rows, left, top, width, height,
              header_bg=None, col_widths=None, alt_row=True,
              header_size=12, body_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # カラム幅の分配
    if col_widths:
        for i, ratio in enumerate(col_widths):
            tbl.columns[i].width = int(width * ratio)
    else:
        cw = width // n_cols
        for i in range(n_cols):
            tbl.columns[i].width = cw

    hbg = header_bg or C_PRIMARY

    # ヘッダー行
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hbg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(header_size)
        run.font.color.rgb = C_WHITE

    # データ行
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if alt_row and ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TABLE_ALT
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(body_size)
            run.font.color.rgb = C_TEXT

    return tbl

def slide_header(slide, title, time_str=None):
    """各スライド共通のヘッダーバー"""
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.65), C_PRIMARY)
    text_box(slide, title,
             Inches(0.3), Inches(0.07), Inches(8.5), Inches(0.55),
             size=18, bold=True, color=C_WHITE)
    if time_str:
        text_box(slide, f"⏱ {time_str}",
                 Inches(7.8), Inches(0.1), Inches(2.0), Inches(0.5),
                 size=12, color=C_PALE_BLUE, align=PP_ALIGN.RIGHT)


# ── スライド定義 ─────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_PRIMARY)
    rect(sl, Inches(0), Inches(3.7), SLIDE_W, Inches(0.06), C_ACCENT)

    text_box(sl, "ざつね屋", Inches(0.6), Inches(0.4), Inches(4), Inches(0.45),
             size=13, color=C_PALE_BLUE)
    text_box(sl, "コース ①", Inches(0.6), Inches(0.95), Inches(9), Inches(0.55),
             size=18, color=C_LIGHT_BLUE)
    text_box(sl, "AI活用知識編", Inches(0.6), Inches(1.5), Inches(9), Inches(1.3),
             size=44, bold=True, color=C_WHITE)
    text_box(sl, "「何から始めればいいかわからない」を解消する 90分",
             Inches(0.6), Inches(2.85), Inches(9), Inches(0.6),
             size=16, color=C_LIGHT_BLUE)
    text_box(sl, "所要時間：90分　｜　成果物：自社 AI活用マップ",
             Inches(0.6), Inches(4.1), Inches(9), Inches(0.5),
             size=13, color=C_PALE_BLUE)


def slide_agenda(prs):
    sl = blank_slide(prs)
    slide_header(sl, "本日のアジェンダ")
    add_table(
        sl,
        ["#", "セクション", "時間"],
        [
            ["1", "AIの種類", "5分"],
            ["2", "できること / できないこと", "5分"],
            ["3", "業務別活用シーン", "10分"],
            ["4", "4軸 AI活用マップ ＋ ツール別表（持ち帰り）", "10分"],
            ["5", "セキュリティと社内ルール", "10分 ＋ 休憩5分"],
            ["6", "プロンプトの基本", "10分"],
            ["7", "実践", "30分"],
            ["8", "まとめ", "5分"],
        ],
        Inches(0.5), Inches(0.8), Inches(9), Inches(4.65),
        col_widths=[0.06, 0.72, 0.22]
    )


def slide_section(prs, num, title, time_str, description=""):
    sl = blank_slide(prs)
    set_bg(sl, C_PRIMARY)

    badge = rect(sl, Inches(0.6), Inches(0.9), Inches(1.2), Inches(1.2), C_ACCENT)
    text_box(sl, str(num),
             Inches(0.6), Inches(0.95), Inches(1.2), Inches(1.1),
             size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    text_box(sl, title,
             Inches(2.2), Inches(1.05), Inches(7.5), Inches(1.0),
             size=30, bold=True, color=C_WHITE)
    text_box(sl, f"⏱ {time_str}",
             Inches(2.2), Inches(2.1), Inches(5), Inches(0.5),
             size=14, color=C_PALE_BLUE)
    if description:
        text_box(sl, description,
                 Inches(2.2), Inches(2.7), Inches(7.5), Inches(1.0),
                 size=15, color=C_LIGHT_BLUE)


def slide_ai_types(prs):
    sl = blank_slide(prs)
    slide_header(sl, "AIの種類")
    add_table(
        sl,
        ["種類", "何をするAIか", "特徴"],
        [
            ["AIエージェント", "指示を受けて自律的にタスクを実行する", "「回答する」だけでなく「実際に動く」"],
            ["動画生成AI", "テキスト・画像から動画を作る", "映像コンテンツを生成・編集する"],
            ["画像生成AI", "テキストから画像を作る", "ビジュアルを生成する"],
            ["会話型AI（チャットAI）", "テキストを理解・生成する", "文章を書く・答える・整理する"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.5),
        col_widths=[0.27, 0.40, 0.33]
    )


def slide_chat_vs_agent(prs):
    sl = blank_slide(prs)
    slide_header(sl, "会話型AI と AIエージェントの違い")
    add_table(
        sl,
        ["", "会話型AI", "AIエージェント"],
        [
            ["できること", "答える・作る・整理する", "実行する・自動化する・連続で動く"],
            ["イメージ", "優秀な相談相手", "自分で動いて仕事を片付ける部下"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.0),
        col_widths=[0.18, 0.41, 0.41]
    )
    rect(sl, Inches(0.5), Inches(3.15), Inches(9), Inches(0.65), C_BG_LIGHT)
    text_box(sl, "💡 このコースで扱うのは主に「会話型AI」です",
             Inches(0.7), Inches(3.22), Inches(8.5), Inches(0.55),
             size=14, bold=True, color=C_PRIMARY)


def slide_can_do(prs):
    sl = blank_slide(prs)
    slide_header(sl, "AIにできること")
    add_table(
        sl,
        ["カテゴリ", "具体例"],
        [
            ["文章を作る", "メール・報告書・SNS投稿・提案書のたたき台"],
            ["文章を整える", "要約・翻訳・校正・言い換え"],
            ["アイデアを出す", "企画案・キャッチコピー・改善案のブレスト"],
            ["情報を整理する", "議事録・調査まとめ・比較表の作成"],
            ["画像・動画を作る", "SNS画像・チラシ・ショート動画"],
            ["業務を自動化する", "ファイル処理・定型作業・ワークフローの自動実行（AIエージェント）"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.55),
        col_widths=[0.27, 0.73]
    )


def slide_cannot_do(prs):
    sl = blank_slide(prs)
    slide_header(sl, "AIにできないこと（注意が必要なこと）")
    add_table(
        sl,
        ["できないこと", "理由"],
        [
            ["最新情報の提供", "学習データに期限がある（ツールによっては検索連携で補える）"],
            ["事実の100%保証", "自信満々に間違える → ハルシネーション（幻覚）と呼ぶ"],
            ["感情・空気を読む", "文脈の裏にある意図は人間が補う必要がある"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.3),
        col_widths=[0.28, 0.72]
    )
    rect(sl, Inches(0.5), Inches(3.35), Inches(9), Inches(0.65), RGBColor(0xFF, 0xF0, 0xE8))
    text_box(sl, "⚠ AIは「下書き・補助ツール」。最終判断は必ず人間が行う",
             Inches(0.7), Inches(3.42), Inches(8.5), Inches(0.55),
             size=14, bold=True, color=C_ACCENT)


def slide_use_cases(prs):
    sl = blank_slide(prs)
    slide_header(sl, "業務別 AI活用シーン")
    add_table(
        sl,
        ["業務", "具体的な活用例", "向いているツール"],
        [
            ["社内外コミュニケーション", "メール文案・議事録・報告書・マニュアル・顧客対応文", "ChatGPT / Gemini / Claude"],
            ["広報・情報発信", "SNS投稿文・チラシ文案・HPテキスト・プレスリリース", "ChatGPT / Gemini / Claude"],
            ["採用", "求人票・会社紹介文・面接質問リスト・内定通知文", "ChatGPT / Gemini / Claude"],
            ["営業・提案", "提案書たたき台・トークスクリプト・お礼メール", "ChatGPT / Gemini / Claude"],
            ["情報収集・調査", "競合調査まとめ・業界トレンド整理・議事録要約", "Gemini / ChatGPT / Claude"],
            ["画像制作", "SNS画像・バナー・チラシ・商品イメージ", "Canva / ChatGPT / Firefly"],
            ["動画・映像", "SNSショート動画・採用動画・会社紹介動画", "CapCut / Canva"],
            ["業務自動化", "CSV整理・ファイル処理・定型作業のスクリプト化", "Claude Code"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.55),
        col_widths=[0.24, 0.46, 0.30]
    )


def slide_ai_map(prs):
    sl = blank_slide(prs)
    slide_header(sl, "4軸 AI活用マップ")

    # 中央の軸線（縦横）
    rect(sl, Inches(4.75), Inches(0.9), Inches(0.04), Inches(4.55), C_BORDER)
    rect(sl, Inches(0.5), Inches(3.1), Inches(9.0), Inches(0.04), C_BORDER)

    # 中央円
    circ = sl.shapes.add_shape(9, Inches(4.35), Inches(2.7), Inches(0.85), Inches(0.85))
    circ.fill.solid()
    circ.fill.fore_color.rgb = C_PRIMARY
    circ.line.fill.background()

    # 4軸ラベル
    text_box(sl, "【文書・テキスト】\nメール／報告書／SNS投稿",
             Inches(3.0), Inches(0.85), Inches(3.5), Inches(0.85),
             size=13, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

    text_box(sl, "【画像・ビジュアル】\nSNS画像／チラシ／バナー",
             Inches(6.5), Inches(2.55), Inches(3.2), Inches(0.85),
             size=13, bold=True, color=RGBColor(0x7C, 0x3A, 0xED))

    text_box(sl, "【動画・映像】\nSNS動画／採用動画",
             Inches(3.0), Inches(4.35), Inches(3.5), Inches(0.85),
             size=13, bold=True, color=C_RED, align=PP_ALIGN.CENTER)

    text_box(sl, "【自動化・開発】\n繰り返し作業の自動化",
             Inches(0.3), Inches(2.55), Inches(3.2), Inches(0.85),
             size=13, bold=True, color=C_ACCENT)

    # 矢印（テキスト代用）
    text_box(sl, "↑", Inches(4.65), Inches(1.7), Inches(0.5), Inches(0.6),
             size=22, color=C_MUTED, align=PP_ALIGN.CENTER)
    text_box(sl, "↓", Inches(4.65), Inches(3.8), Inches(0.5), Inches(0.6),
             size=22, color=C_MUTED, align=PP_ALIGN.CENTER)
    text_box(sl, "←", Inches(3.3), Inches(2.9), Inches(1.0), Inches(0.5),
             size=22, color=C_MUTED, align=PP_ALIGN.RIGHT)
    text_box(sl, "→", Inches(5.2), Inches(2.9), Inches(1.0), Inches(0.5),
             size=22, color=C_MUTED)


def slide_ai_map_table(prs):
    sl = blank_slide(prs)
    slide_header(sl, "4軸 AI活用マップ — 領域別概要")
    add_table(
        sl,
        ["領域", "対象スタッフ", "主な用途", "課金"],
        [
            ["文書・テキスト", "全スタッフ", "メール・報告書・広報文", "無料で可"],
            ["画像・ビジュアル", "広報・営業・SNS担当", "SNS画像・チラシ・バナー", "無料で可"],
            ["動画・映像", "広報・採用・SNS担当", "SNS動画・採用動画", "無料で可"],
            ["自動化・開発", "経営者・業務改善担当", "繰り返し作業の自動化", "APIキー必要"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.8),
        col_widths=[0.22, 0.25, 0.31, 0.22]
    )
    text_box(sl, "→ 次のページ：ツール別詳細（持ち帰り資料）",
             Inches(0.5), Inches(3.85), Inches(9), Inches(0.5),
             size=13, italic=True, color=C_MUTED)


def slide_tool_table1(prs):
    sl = blank_slide(prs)
    slide_header(sl, "ツール別詳細 ① 文書系（持ち帰り資料）")
    add_table(
        sl,
        ["ツール", "手間", "学習コスト", "専門知識", "課金"],
        [
            ["ChatGPT", "★☆☆", "★☆☆", "★☆☆", "無料あり／Plus 約3,000円/月"],
            ["Gemini", "★☆☆", "★☆☆", "★☆☆", "無料あり／Advanced 2,900円/月"],
            ["Microsoft Copilot", "★★☆", "★☆☆", "★☆☆", "M365契約前提／Copilot 約4,500円/ユーザー/月"],
            ["Claude", "★☆☆", "★☆☆", "★☆☆", "無料あり／Pro 約3,000円/月"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.7),
        col_widths=[0.22, 0.10, 0.13, 0.13, 0.42]
    )
    note = sl.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.6))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "手間：★☆☆＝アカウントのみ　★★☆＝インストール・設定あり　★★★＝環境構築が必要"
    r.font.size = Pt(10)
    r.font.color.rgb = C_MUTED
    add_para(tf, "学習：★☆☆＝数時間で実務に使える　★★☆＝数日の練習が必要　★★★＝継続的な学習が必要",
             size=10, color=C_MUTED)
    add_para(tf, "専門知識：★☆☆＝ITスキル不問　★★☆＝基本的なPC操作が必要　★★★＝プログラミングの概念理解が必要",
             size=10, color=C_MUTED)


def slide_tool_table2(prs):
    sl = blank_slide(prs)
    slide_header(sl, "ツール別詳細 ② 画像・動画・自動化（持ち帰り資料）")
    add_table(
        sl,
        ["ツール", "領域", "手間", "学習コスト", "課金"],
        [
            ["Canva（AI機能）", "画像", "★☆☆", "★☆☆", "無料あり／Pro 1,500円/月"],
            ["Bing Image Creator", "画像", "★☆☆", "★☆☆", "完全無料"],
            ["Adobe Firefly", "画像", "★☆☆", "★★☆", "無料あり（25クレジット/月）／CC 6,480円/月〜"],
            ["CapCut（AI機能）", "動画", "★★☆", "★★☆", "無料あり／Pro 約1,200円/月"],
            ["Canva（動画）", "動画", "★☆☆", "★☆☆", "無料あり／Pro 1,500円/月"],
            ["Claude Code", "自動化", "★★★", "★★★", "APIキー従量課金（軽い使い方で月数百円〜）"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.0),
        col_widths=[0.23, 0.10, 0.10, 0.13, 0.44]
    )


def slide_security_input(prs):
    sl = blank_slide(prs)
    slide_header(sl, "入力してはいけない情報")
    add_table(
        sl,
        ["カテゴリ", "具体例"],
        [
            ["個人情報", "顧客名・住所・電話番号・メールアドレス"],
            ["機密情報", "未発表の事業計画・契約内容・取引先との交渉内容"],
            ["社内情報", "従業員の給与・人事評価・内部トラブル"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(2.2),
        col_widths=[0.22, 0.78]
    )
    rect(sl, Inches(0.5), Inches(3.25), Inches(9), Inches(0.65), RGBColor(0xFF, 0xF0, 0xE8))
    text_box(sl, "❗ 入力した情報がAIの学習に使われる可能性があります（ツールにより異なる）",
             Inches(0.7), Inches(3.32), Inches(8.5), Inches(0.55),
             size=14, bold=True, color=C_ACCENT)


def slide_security_learning(prs):
    sl = blank_slide(prs)
    slide_header(sl, "ツール別 学習データ利用状況")
    add_table(
        sl,
        ["ツール", "学習利用", "備考"],
        [
            ["ChatGPT", "デフォルトON", "設定でOFFにできる"],
            ["Gemini", "デフォルトON", "設定でOFFにできる"],
            ["Claude", "デフォルトOFF", "商用利用は学習に使わない"],
            ["Microsoft Copilot", "OFF", "M365契約内は学習利用なし"],
            ["Claude Code", "OFF", "API経由のため学習利用なし・社内情報を扱う業務自動化に向いている"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.3),
        col_widths=[0.24, 0.18, 0.58]
    )
    rect(sl, Inches(0.5), Inches(4.35), Inches(9), Inches(0.65), C_BG_LIGHT)
    text_box(sl, "💡 社内の機密情報を扱う業務には、学習OFFのツールを選ぶのが安心です",
             Inches(0.7), Inches(4.42), Inches(8.5), Inches(0.55),
             size=14, bold=True, color=C_PRIMARY)


def slide_security_rules(prs):
    sl = blank_slide(prs)
    slide_header(sl, "社内ルールの作り方（最低限）＆ 出力を使う際の注意")

    # 左：社内ルール
    text_box(sl, "社内ルール（最低限）",
             Inches(0.5), Inches(0.82), Inches(4.2), Inches(0.5),
             size=14, bold=True, color=C_PRIMARY)

    rules = sl.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.2), Inches(2.8))
    tf = rules.text_frame
    tf.word_wrap = True
    items = [
        "① 使ってよいツールを決める",
        "② 入力禁止情報のリストを共有する",
        "③ AIの出力は必ず確認してから使う",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = item
        r.font.size = Pt(14)
        r.font.color.rgb = C_TEXT
        if i < len(items) - 1:
            gap = tf.add_paragraph()
            gap.add_run().text = ""

    # 区切り線
    rect(sl, Inches(4.95), Inches(0.82), Inches(0.04), Inches(4.0), C_BORDER)

    # 右：出力の注意
    text_box(sl, "出力を使うときの注意",
             Inches(5.2), Inches(0.82), Inches(4.5), Inches(0.5),
             size=14, bold=True, color=C_ACCENT)

    cautions = sl.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.5), Inches(3.5))
    tf2 = cautions.text_frame
    tf2.word_wrap = True
    clist = [
        ("ファクトチェック必須",
         "→ AIは自信満々に嘘をつく（ハルシネーション）"),
        ("著作権の確認",
         "→ 生成画像・文章の商用利用は各ツールの規約を確認"),
        ("最終判断は人間が行う",
         "→ AIはあくまで下書き・補助"),
    ]
    for i, (head, sub) in enumerate(clist):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        r = p.add_run()
        r.text = head
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = C_TEXT
        add_para(tf2, sub, size=12, color=C_MUTED)
        if i < len(clist) - 1:
            tf2.add_paragraph().add_run().text = ""


def slide_prompt_basics(prs):
    sl = blank_slide(prs)
    slide_header(sl, "良いプロンプトの4要素")
    add_table(
        sl,
        ["要素", "内容", "例"],
        [
            ["役割", "AIにどの立場で考えるかを伝える", "「あなたは中小企業向けの採用担当者です」"],
            ["背景", "状況・前提を伝える", "「広島県の製造業・従業員30名」"],
            ["指示", "何をしてほしいか具体的に", "「求人票を作成してください」"],
            ["形式", "出力の形を指定する", "「箇条書きで・400字以内で」"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.0),
        col_widths=[0.12, 0.38, 0.50]
    )
    rect(sl, Inches(0.5), Inches(4.05), Inches(9), Inches(0.65), C_BG_LIGHT)
    text_box(sl, "💡 4要素がすべて揃うほど、期待通りの出力が得られやすくなります",
             Inches(0.7), Inches(4.12), Inches(8.5), Inches(0.55),
             size=14, bold=True, color=C_PRIMARY)


def slide_before_after(prs):
    sl = blank_slide(prs)
    slide_header(sl, "Before / After — プロンプト改善例")

    text_box(sl, "❌ Before（悪い例）",
             Inches(0.5), Inches(0.82), Inches(9), Inches(0.4),
             size=13, bold=True, color=C_RED)
    rect(sl, Inches(0.5), Inches(1.22), Inches(9), Inches(0.72), RGBColor(0xFF, 0xF0, 0xF0))
    text_box(sl, "「求人票を書いて」",
             Inches(0.7), Inches(1.3), Inches(8.5), Inches(0.6),
             size=16, bold=True, color=C_RED)

    text_box(sl, "↓  4要素を追加する",
             Inches(0.5), Inches(2.05), Inches(9), Inches(0.45),
             size=13, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    text_box(sl, "✅ After（良い例）",
             Inches(0.5), Inches(2.6), Inches(9), Inches(0.4),
             size=13, bold=True, color=C_GREEN)
    rect(sl, Inches(0.5), Inches(3.0), Inches(9), Inches(1.6), RGBColor(0xEC, 0xFD, 0xF5))
    after = ("「あなたは中小企業向けの採用担当者です。\n"
             "広島県の製造業・従業員30名の会社で、未経験歓迎の溶接工を募集しています。\n"
             "応募が集まりやすい求人票を箇条書きで作成してください」")
    text_box(sl, after,
             Inches(0.7), Inches(3.08), Inches(8.5), Inches(1.45),
             size=14, color=RGBColor(0x06, 0x5F, 0x46))


def slide_hallucination(prs):
    sl = blank_slide(prs)
    slide_header(sl, "ハルシネーション（AIの嘘）を減らす指示の出し方")
    add_table(
        sl,
        ["対策", "プロンプトへの追記方法"],
        [
            ["わからないことを認めさせる", "「不明な点は『わかりません』と答えてください」をプロンプトに加える"],
            ["確信度を明示させる", "「確信が持てない情報には【要確認】とつけてください」と指示する"],
            ["推測を排除する", "「推測や不確かな情報は含めないでください」と指示する"],
            ["情報はこちらが提供する", "「以下の情報をもとに整理してください」と事実を渡してから指示する"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(3.0),
        col_widths=[0.30, 0.70]
    )
    rect(sl, Inches(0.5), Inches(4.05), Inches(9), Inches(0.65), C_BG_LIGHT)
    text_box(sl, "💡 ハルシネーションはAIの限界ではなく、指示の出し方で大幅に減らせます",
             Inches(0.7), Inches(4.12), Inches(8.5), Inches(0.55),
             size=14, bold=True, color=C_PRIMARY)


def slide_practice(prs):
    sl = blank_slide(prs)
    slide_header(sl, "実践（30分） — ツール：ChatGPT または Gemini（受講者アカウント優先）")
    add_table(
        sl,
        ["Step", "内容", "時間"],
        [
            ["1", "お題を選ぶ（社内向けお知らせ ／ SNS投稿文 ／ 求人票のたたき台）", "2分"],
            ["2", "4要素を使ってプロンプトを自分で書く", "5分"],
            ["3", "AIに投げて結果を見る", "3分"],
            ["4", "修正指示を出してみる（「もっと短く」「やわらかく」等）", "5分"],
            ["5", "全体シェア（1〜2人が披露）", "5分"],
            ["6", "意見交換（受講者同士でフィードバック）", "5分"],
            ["", "⭐ バッファ・Q&A", "5分"],
        ],
        Inches(0.5), Inches(0.85), Inches(9), Inches(4.35),
        col_widths=[0.07, 0.73, 0.20]
    )


def slide_summary(prs):
    sl = blank_slide(prs)
    slide_header(sl, "まとめ（5分）")

    text_box(sl, "本日の学び",
             Inches(0.5), Inches(0.82), Inches(5.5), Inches(0.45),
             size=14, bold=True, color=C_PRIMARY)

    learns = sl.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.2))
    tf = learns.text_frame
    tf.word_wrap = True
    for i, item in enumerate([
        "✅ AIにできること・できないことがわかった",
        "✅ 自社業務にどのツールが合うかがわかった",
        "✅ 安心して使える範囲（セキュリティ）がわかった",
        "✅ プロンプトの出し方の基本がわかった",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = item
        r.font.size = Pt(13)
        r.font.color.rgb = C_TEXT

    text_box(sl, "持ち帰り資料",
             Inches(6.3), Inches(0.82), Inches(3.4), Inches(0.45),
             size=14, bold=True, color=C_ACCENT)

    resources = sl.shapes.add_textbox(Inches(6.3), Inches(1.3), Inches(3.4), Inches(1.6))
    tf2 = resources.text_frame
    tf2.word_wrap = True
    for i, item in enumerate([
        "📋 ツール別詳細一覧（別表）",
        "📋 プロンプト4要素メモ",
        "📋 社内ルール作成チェックリスト",
    ]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        r = p.add_run()
        r.text = item
        r.font.size = Pt(13)
        r.font.color.rgb = C_TEXT

    # サービス案内バナー
    rect(sl, Inches(0.5), Inches(3.6), Inches(9), Inches(1.8), C_BG_LIGHT)
    text_box(sl, "次のステップ | ざつね屋のサービス案内",
             Inches(0.7), Inches(3.7), Inches(8.5), Inches(0.45),
             size=13, bold=True, color=C_PRIMARY)
    svc = sl.shapes.add_textbox(Inches(0.7), Inches(4.18), Inches(8.5), Inches(1.1))
    tf3 = svc.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "📚 実践編（文書・画像・動画・Claude Code）— 今日の知識を体験に変える"
    r3.font.size = Pt(12)
    r3.font.color.rgb = C_TEXT
    add_para(tf3, "🛠 AI業務改善オーダーメイド・🤝 AI開発伴走 — 自社の課題を一緒に解決する",
             size=12, color=C_TEXT)


# ── メイン ──────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)
    slide_agenda(prs)

    slide_section(prs, "1", "AIの種類", "5分", "AIにはどんな種類があるか、全体を掴みましょう")
    slide_ai_types(prs)
    slide_chat_vs_agent(prs)

    slide_section(prs, "2", "できること / できないこと", "5分", "期待値を正しく調整することが、AI活用の第一歩です")
    slide_can_do(prs)
    slide_cannot_do(prs)

    slide_section(prs, "3", "業務別活用シーン", "10分", "「自分の仕事に使えそう」を自分ごと化する")
    slide_use_cases(prs)

    slide_section(prs, "4", "4軸 AI活用マップ", "10分", "AIをどの方向に使うかを地図で整理する")
    slide_ai_map(prs)
    slide_ai_map_table(prs)
    slide_tool_table1(prs)
    slide_tool_table2(prs)

    slide_section(prs, "5", "セキュリティと社内ルール", "10分 ＋ 休憩5分", "安心して使えるルールを確認しましょう")
    slide_security_input(prs)
    slide_security_learning(prs)
    slide_security_rules(prs)

    slide_section(prs, "6", "プロンプトの基本", "10分", "指示の質がそのまま出力の質になります")
    slide_prompt_basics(prs)
    slide_before_after(prs)
    slide_hallucination(prs)

    slide_section(prs, "7", "実践", "30分", "体験して「できる」を感じて帰りましょう")
    slide_practice(prs)

    slide_section(prs, "8", "まとめ", "5分", "本日の学びとサービス案内")
    slide_summary(prs)

    output = "training-materials/course01_AI活用知識編.pptx"
    prs.save(output)
    print(f"[OK] 生成完了: {output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
