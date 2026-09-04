# -*- coding: utf-8 -*-
"""サービス説明版スライド（比較資料の凝縮版・13枚・投影向けライトテーマ）。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUT = r"C:\Users\ooto\work\ClaudeCode\kurashi-no-dodai-log\00-01_han-ai\docs\業務資料\01_入口_営業と見極め\サービス説明版スライド.pptx"
INK = RGBColor(0x1A, 0x22, 0x33)
SUB = RGBColor(0x5B, 0x64, 0x72)
ACC = RGBColor(0x2F, 0x6F, 0x5E)
ACCL = RGBColor(0xE4, 0xEF, 0xEA)
PANEL = RGBColor(0xF2, 0xF7, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
KICK = RGBColor(0xD8, 0xE8, 0xE2)
JP = "Yu Gothic"
ML = 0.9
CW = 13.333 - 2 * ML

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def sf(run, size, bold=False, color=INK, name=JP):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def nsh(s):
    s.shadow.inherit = False


def content(title):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(ML), Inches(0.42), Inches(CW), Inches(0.8))
    sf(tb.text_frame.paragraphs[0].add_run(), 25, True, INK)
    tb.text_frame.paragraphs[0].runs[0].text = title
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(1.2), Inches(2.6), Pt(3))
    r.fill.solid(); r.fill.fore_color.rgb = ACC; r.line.fill.background(); nsh(r)
    n = len(prs.slides._sldIdLst)
    ft = s.shapes.add_textbox(Inches(9.6), Inches(7.04), Inches(2.83), Inches(0.32))
    p = ft.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    rr = p.add_run(); rr.text = "ざつね屋　｜　{}".format(n); sf(rr, 8, color=SUB)
    return s


def lines(s, items, top=1.55, width=CW, left=ML, box_h=None):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                              Inches(box_h if box_h else 7.5 - top - 0.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        t, sz, b, c, sa = (it + (None,) * 5)[:5]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(sa if sa is not None else 6)
        rn = p.add_run(); rn.text = t
        sf(rn, sz or 15, b or False, c or INK)
    return tb


def table(s, rows, top, col_w, fs=12, row_h=0.5, first_bold=True):
    nr, nc = len(rows), len(rows[0])
    gf = s.shapes.add_table(nr, nc, Inches(ML), Inches(top), Inches(sum(col_w)), Inches(row_h * nr))
    t = gf.table
    t.first_row = False; t.horz_banding = False
    for ci, w in enumerate(col_w):
        t.columns[ci].width = Inches(w)
    for ri, row in enumerate(rows):
        t.rows[ri].height = Inches(row_h)
        head = ri == 0
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.09); cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCL if head else WHITE
            tf = cell.text_frame; tf.word_wrap = True; tf.clear()
            for li, ln in enumerate(str(val).split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                rn = p.add_run(); rn.text = ln
                sf(rn, fs, head or (first_bold and ci == 0), INK)
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                ln = tcPr.makeelement(qn(tag), {"w": "9525", "cap": "flat", "cmpd": "sng", "algn": "ctr"})
                fill = ln.makeelement(qn("a:solidFill"), {})
                clr = fill.makeelement(qn("a:srgbClr"), {"val": "C9D2CE"})
                fill.append(clr); ln.append(fill); tcPr.append(ln)
    return gf


def callout(s, text, top, height=0.95):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(top), Inches(0.07), Inches(height))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACC; bar.line.fill.background(); nsh(bar)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML + 0.07), Inches(top), Inches(CW - 0.07), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = PANEL; box.line.fill.background(); nsh(box)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.16)
    sf(tf.paragraphs[0].add_run(), 12.5, color=INK)
    tf.paragraphs[0].runs[0].text = text


def notes(s, t):
    s.notes_slide.notes_text_frame.text = t


# 1 表紙
s = prs.slides.add_slide(BLANK)
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(2.4), Inches(2.6), Pt(4))
r.fill.solid(); r.fill.fore_color.rgb = ACC; r.line.fill.background(); nsh(r)
tb = s.shapes.add_textbox(Inches(ML), Inches(2.65), Inches(CW), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; rn = p.add_run(); rn.text = "業務改善・BPR・DX"; sf(rn, 40, True, INK)
p = tf.add_paragraph(); rn = p.add_run(); rn.text = "── 御社は、どこから手をつけると成果が出やすいか"; sf(rn, 20, False, SUB)
tb2 = s.shapes.add_textbox(Inches(ML), Inches(6.3), Inches(CW), Inches(0.6))
rn = tb2.text_frame.paragraphs[0].add_run(); rn.text = "ざつね屋　│　地域企業の業務変革のご相談"; sf(rn, 12, color=SUB)
notes(s, "3つは優劣ではなく変える範囲が違うだけ。今どの段階かで着手点が決まる、という話をする。")

# 2 立場
s = content("この資料の立場")
lines(s, [
    ("「業務改善・BPR・DX、どれをやるべきか」に、順位では答えません。3つは変える範囲が違うだけです。", 16, False, INK, 12),
    ("お答えするのは1点 ── 御社の今の状態なら、どこから手をつけると一番早く成果が出るか。", 17, True, ACC, 12),
    ("デジタルもAIも、そのための手段です。目的が先、手段は後。", 15, False, INK, 6),
], top=1.7)
callout(s, "ざつね屋の立場：成果を売る。デジタルとAIは、そのための手段にすぎない。", 5.4)

# 3 3つの定義
s = content("3つのアプローチ")
table(s, [
    ["アプローチ", "変える範囲", "ひとことで言うと"],
    ["業務改善設計", "個人・部門内の作業", "今のやり方を前提に、ムダな作業を減らす"],
    ["BPR設計", "部門をまたぐ業務プロセス", "流れをゼロから作り直す（まず徹底的に見直す）"],
    ["DX設計", "事業・組織・企業文化", "データとデジタルで、事業そのものを変える"],
], top=1.7, col_w=[2.4, 3.6, 5.53], fs=13, row_h=1.0)
callout(s, "外側は内側を含む：DXを進めれば必ずプロセスの作り直し（BPR）が起き、その中で作業の効率化（業務改善）が起きる。", 6.0, 0.9)

# 4 包含図
s = content("3つの関係")
o = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.667), Inches(1.55), Inches(8.0), Inches(4.4))
o.fill.background(); o.line.color.rgb = ACC; o.line.width = Pt(2.5); nsh(o)
m = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.667), Inches(2.3), Inches(6.0), Inches(2.95))
m.fill.solid(); m.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xF6); m.line.color.rgb = ACC; m.line.width = Pt(1.75); nsh(m)
i = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.667), Inches(3.0), Inches(4.0), Inches(1.6))
i.fill.solid(); i.fill.fore_color.rgb = ACCL; i.line.color.rgb = ACC; i.line.width = Pt(1.5); nsh(i)
for shp, name, cap in [(o, "DX設計", "経営インパクト大・着手しやすさ小"),
                       (m, "BPR設計", "インパクト中・着手しやすさ中"),
                       (i, "業務改善設計", "インパクト小・着手しやすさ大")]:
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP if shp is not i else MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rn = p.add_run(); rn.text = name; sf(rn, 15, True, ACC)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    rn = p2.add_run(); rn.text = cap; sf(rn, 9.5, False, SUB)
tb = s.shapes.add_textbox(Inches(ML), Inches(6.2), Inches(CW), Inches(0.7))
rn = tb.text_frame.paragraphs[0].add_run()
rn.text = "内側から外側へ育てられる：業務改善で足場をつくり → BPRで作り直し → DXで事業に効かせる。中小企業ではこの順が現実的。"
sf(rn, 11, color=INK)

# 5 早見表
s = content("早見表：主な問い")
table(s, [
    ["アプローチ", "主な問い", "経営との距離"],
    ["業務改善設計", "この作業は、もっと楽にできないか", "遠い（現場の判断で始められる）"],
    ["BPR設計", "この流れ自体、そもそも必要か。作り直すとどうなるか", "中間（部門長・経営者の判断が要る）"],
    ["DX設計", "この事業は、これからも同じやり方で続けられるのか", "近い（経営者の意思決定そのもの）"],
], top=1.7, col_w=[2.3, 6.2, 3.03], fs=12.5, row_h=1.0)
callout(s, "ご相談の内容が、どの「主な問い」に近いか ── そこで入口の見当がつきます。", 6.0, 0.85)

# 6-8 各アプローチ
def approach(title, lead, flow, cases):
    s = content(title)
    lines(s, [(lead, 13, False, SUB, 10)], top=1.5, box_h=0.9)
    lines(s, [("進め方：" + flow, 13, False, INK, 12)], top=2.4, box_h=1.4)
    lines(s, [("向いているのは", 14, True, ACC, 4)] + [("・" + c, 13, False, INK, 4) for c in cases], top=3.9)
    return s

approach("① 業務改善設計",
         "今の業務のやり方を前提に、その一部を効率化します。対象は「作業」。成果は「時間が減る／ミスが減る／担当者が楽になる」。着手しやすく、失敗時の損失も小さい。",
         "対象業務を1つ選ぶ → 作業を洗い出す → 現状を数字で見る → ECRS（なくす・まとめる・入れ替える・簡単にする）で改善案 → 小さく試す → 手順書にして定着",
         ["業務の全体像をまだ誰も書き出していない",
          "特定の担当者に負担が偏っている",
          "大きな投資判断の前に、小さな成功例が必要"])
approach("② BPR設計",
         "業務プロセスそのものを抜本的に見直し、ゼロベースで作り直します。対象は「作業」ではなく「流れ」。今のやり方を前提にしません。",
         "業務分析（現状把握 → 問題の特定 → 課題設定）→ 施策検討（A/B/C案を比較 → 実施計画）→ 施策実施（KPIを決めて実行 → 効果検証）",
         ["流れ自体に無理がある（同じ入力の繰り返し、承認が何段も）",
          "部門をまたぐと途端に止まる",
          "システムを入れたのに紙の作業も残っている（二重作業）"])
approach("③ DX設計",
         "データとデジタル技術で、事業・組織・企業文化そのものを変革し、経営としての成果を生みます。その中で必要なプロセスの作り直し（BPR）も行います。",
         "経営ビジョンを決める → 現状とデータ資産の棚卸し → 変革テーマを絞る → 小さく試す（PoC）→ 仕組み化・スケール → 組織・文化の定着",
         ["今のやり方のままでは続かない、という危機感が経営者にある",
          "データが蓄積されているが活かせていない",
          "事業モデル・組織・評価の仕方まで変える覚悟がある"])

# 9 決定木
s = content("どこから着手するか")
qx, qw, ox, ow = 2.3, 4.7, 8.1, 4.3
rows = [
    ("経営層の関与・予算/体制はあるか", 1.7, ("弱い・未整備", "① 業務改善設計（小さな成功体験を作る）")),
    ("業務の全体像を書き出せているか", 2.75, ("できていない", "① 業務改善設計（洗い出しから）")),
    ("流れ自体に無理／部門をまたぐと止まるか", 3.8, ("はい", "② BPR設計（プロセスを作り直す）")),
    ("データ活用・事業モデルの変革まで踏み込むか", 4.95, ("はい", "③ DX設計（中でBPRを実施）")),
]
prevb = None
for txt, y, (bl, ot) in rows:
    h = 0.8
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(qx), Inches(y), Inches(qw), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = ACC; box.line.width = Pt(1.4); nsh(box)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rn = p.add_run(); rn.text = txt; sf(rn, 11, False, INK)
    if prevb is not None:
        cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(qx + qw / 2), Inches(prevb), Inches(qx + qw / 2), Inches(y))
        cn.line.color.rgb = ACC; cn.line.width = Pt(1.2); nsh(cn)
    prevb = y + h
    mid = y + h / 2
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(qx + qw + 0.02), Inches(mid), Inches(ox - 0.02), Inches(mid))
    cn.line.color.rgb = ACC; cn.line.width = Pt(1.2); nsh(cn)
    obx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ox), Inches(y), Inches(ow), Inches(h))
    obx.fill.solid(); obx.fill.fore_color.rgb = ACCL; obx.line.color.rgb = ACC; obx.line.width = Pt(1.2); nsh(obx)
    tf = obx.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rn = p.add_run(); rn.text = ot; sf(rn, 10, False, INK)
    lb = s.shapes.add_textbox(Inches(qx + qw + 0.05), Inches(mid - 0.32), Inches(1.0), Inches(0.3))
    rn = lb.text_frame.paragraphs[0].add_run(); rn.text = bl; sf(rn, 8.5, color=SUB)

# 10 判断表
s = content("ご相談内容からの逆引き")
table(s, [
    ["こういう状態なら", "着手するアプローチ"],
    ["業務の全体像を誰も書き出していない", "① 業務改善設計（洗い出しから）"],
    ["同じ入力の繰り返し／承認が何段も／部門をまたぐと止まる", "② BPR設計"],
    ["システムを入れたが紙の作業も残っている", "② BPR設計"],
    ["データを活かしたい／事業モデル・組織・文化まで変えたい", "③ DX設計（中でBPR）"],
    ["経営層の関与が弱い／予算・体制が未整備", "① 業務改善設計（土台づくり）"],
], top=1.6, col_w=[7.4, 4.13], fs=11.5, row_h=0.72)

# 11 進め方
s = content("ざつね屋の進め方")
lines(s, [
    ("1. 入口 ── 初回ヒアリングと簡易診断で、どのアプローチから入るかを見極める", 14, False, INK, 8),
    ("2. アプローチ別に伴走 ── 業務改善／BPR／DX それぞれの手順と道具（記入シート）で進める", 14, False, INK, 8),
    ("3. 成果を数字にする ── どのアプローチでも、着手前の数値を測り、KPIで効果を確認する", 14, False, INK, 8),
    ("4. 自走へ ── 手順書・チェックリスト・進め方ガイドで、御社だけで回せる状態にする", 14, False, INK, 8),
], top=1.7)
callout(s, "「まず1つの業務／1つの段階」から始められる、段階契約に対応しています。", 5.2, 0.85)

# 12 体制
s = content("体制と関わり方")
table(s, [
    ["役割", "担当"],
    ["進行・設計・分析・現場伴走・研修", "ざつね屋"],
    ["業務の説明・作業記録の協力", "御社の現場担当"],
    ["改善案・施策の意思決定", "御社の責任者"],
], top=1.7, col_w=[6.6, 4.93], fs=13, row_h=0.75)
lines(s, [
    ("含まないもの：外部システムの構築・保守、データ入力の代行", 12, False, SUB, 4),
    ("段階の区切りごとに結果を見て、次に進むかを判断できます。", 12, False, SUB, 4),
], top=4.6)

# 13 まとめ
s = content("まとめ")
lines(s, [
    ("3つは排他的なものではありません。", 15, False, INK, 10),
    ("業務改善で足場を作り → BPRで作り直し → DXで事業に効かせる。", 18, True, ACC, 10),
    ("今どの段階にいるかを見誤らないことが、遠回りを避ける一番の近道です。", 15, False, INK, 16),
    ("まずは初回ヒアリングから。御社の業務を一緒に見立てます。", 14, False, INK, 6),
    ("お問い合わせ：ざつね屋　zatuneya@gmail.com", 13, True, INK, 6),
], top=1.7)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
